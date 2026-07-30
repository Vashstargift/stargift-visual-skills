# -*- coding: utf-8 -*-
"""Сборщик PowerPoint-версии премиальных презентаций Stargift (эталон v2).

Повторяет раскладки слайдов из constructor.html (функции TPL / trioTPL) в .pptx:
те же цвета, отступы, пропорции и кегли для каждого типа слайда. Текст —
редактируемые текст-боксы (менеджер правит в PowerPoint / Keynote), фото лотов
вписываются с сохранением пропорций (contain), фоновые фото cover/divider —
кроп «cover» + полупрозрачное затемнение.

ВХОД:
    build_pptx(data, out_path)
    data = {"title": str, "slides": [ {"type": ..., "params": {...}, "lots":[...]} ]}

Зависимости: только python-pptx и stdlib (PIL — опционально, для .webp).
Холст 16:9, 1920×1080 px. 1 px = 6350 EMU (914400/144); кегль в pt = px * 0.5.
"""
import base64
import re
import tempfile
import urllib.request
import urllib.parse
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── единицы: холст 1920×1080 px → слайд 13.333×7.5" ──
PXE = 914400 / 144.0            # EMU в одном px (= 6350)
def E(px):  return Emu(int(round(px * PXE)))     # px → EMU
def PT(px): return Pt(px * 0.5)                  # px-кегль → pt (960pt на 1920px)
def C(h):   return RGBColor.from_string(h.lstrip("#").upper())

# ── палитра эталона ──
DARK      = "000000"   # тёмный фон
PAPER     = "FFFFFF"   # бумага
CARD_L    = "F5F5F5"   # светлая карточка
CARD_L_BD = "E2E2E2"   # бордер светлой карточки
CARD_D    = "1C1C1C"   # тёмная карточка
GOLD      = "B9975B"   # золото (кикер на тёмном)
GOLD_HL   = "D4B87E"   # золото цены на тёмном
GOLD_L    = "96702F"   # золото на светлом
MUTED     = "8A8A8A"   # приглушённый кикер/сертификатор
TXT_L     = "F5F5F5"   # светлый текст
TXT_L2    = "D9D9D9"   # светлый текст (второй)
TXT_D     = "111111"   # тёмный текст
TXT_D2    = "3A3A3A"   # тёмный текст (второй)
DESC_D    = "B8B8B8"   # описание на тёмном
DESC_L    = "4D4D4D"   # описание на светлом
LINE_D    = "303030"   # линия на тёмном
LINE_L    = "DCDCDC"   # линия на светлом
LINE_L2   = "DCDCDC"   # линия-разделитель списков

PRATA = "Prata"        # заголовки/цены (serif)
GOLOS = "Golos Text"   # основной текст

# Логотипы (локальные файлы на машине воркера; нет файла — фолбэк на текст «STARGIFT»)
_IMGDIR = Path(__file__).resolve().parent.parent / "img"
LOGO_WHITE = _IMGDIR / "stargift_logo_white.png"   # светлая версия — на тёмном фоне
LOGO_DARK = _IMGDIR / "stargift_logo.png"          # тёмная версия — на светлом фоне

NBSP = " "


def _fmt_num(n):
    """1500000 → «1 500 000» (неразрывные пробелы)."""
    s = str(int(round(float(n))))
    out = ""
    for i, ch in enumerate(reversed(s)):
        if i and i % 3 == 0:
            out = NBSP + out
        out = ch + out
    return out


def _fmt_rub(n):
    return _fmt_num(n) + NBSP + "₽"


def _plural(n, one, few, many):
    m10, m100 = n % 10, n % 100
    if m10 == 1 and m100 != 11:
        return one
    if 2 <= m10 <= 4 and not (10 <= m100 < 20):
        return few
    return many


def _roman(n):
    m = ["", "I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X",
         "XI", "XII", "XIII", "XIV", "XV", "XVI"]
    return m[n] if 0 <= n < len(m) else str(n)


_TAG_RE = re.compile(r"<[^>]+>")


def _fit_sent(txt, max_chars):
    """Обрезка по последнему ЦЕЛОМУ предложению в бюджете (правка Вашика 30.07:
    «…» посреди фразы запрещены — предложения заканчиваются всегда)."""
    txt = (txt or "").strip()
    if len(txt) <= max_chars:
        return txt
    cut = txt[:max_chars]
    best = -1
    for i in range(len(cut) - 1, -1, -1):
        ch = cut[i]
        if ch in ".!?":
            best = i; break
        if ch == "»" and i > 0 and cut[i-1] in ".!?":
            best = i; break
    if best > 20:
        return cut[:best + 1]
    import re as _re2
    m = _re2.match(r"^[^.!?]*[.!?]»?", txt)
    return m.group(0) if m else txt


def _clean(s):
    """Убрать HTML-разметку из текста лота (в resolved-данных имена/описания
    иногда приходят с тегами <p …>) и схлопнуть пробелы."""
    if not s:
        return ""
    return " ".join(_TAG_RE.sub(" ", str(s)).split())


def _all_desc_empty(lots):
    """True, если у ВСЕХ лотов слайда описание пустое/пробельное.
    В этом режиме фото-бокс делаем выше, а текстовую зону — компактнее."""
    lots = lots or []
    return bool(lots) and all(not (l.get("desc") or "").strip() for l in lots)


def _avail_line(lot):
    """Наличие без точек-разделителей и галерей (правка Вашика 28.07)."""
    return "В наличии в Москве" if lot.get("instock") else "Под заказ, доставка около 2 месяцев"


def _spec_line(lot):
    """Спецификация: размер, сертификатор (через запятую, без «Provenance»)."""
    parts = []
    if lot.get("size"):
        parts.append(lot["size"])
    cert = (lot.get("cert") or "").strip()
    if cert and cert != "Provenance":
        parts.append(cert)
    return ", ".join(parts)


def _meta_line(lot):
    parts = [p for p in (_spec_line(lot), _avail_line(lot)) if p]
    return ", ".join(parts)


# ──────────────────────────────────────────────────────────────────────────
#  Загрузка фото (https-URL или data:-URI) во временные файлы
# ──────────────────────────────────────────────────────────────────────────
_UA = ("Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 "
       "(KHTML, like Gecko) Chrome/120.0 Safari/537.36")


def _ext_of(url):
    low = url.lower().split("?")[0]
    for e in (".png", ".jpeg", ".jpg", ".gif", ".bmp", ".tiff", ".webp"):
        if low.endswith(e):
            return e
    return ".jpg"


def _to_supported(path, workdir, idx):
    """python-pptx не понимает .webp — конвертируем через PIL, если доступен."""
    if path.suffix.lower() != ".webp":
        return path
    try:
        from PIL import Image
        dst = workdir / f"cv_{idx}.png"
        Image.open(path).convert("RGB").save(dst, "PNG")
        return dst
    except Exception:
        return None


_FETCH_CACHE = {}


def _fetch(src, workdir, idx):
    """Скачать/декодировать фото. None при неудаче — слайд собирается без него.
    Один и тот же URL в пределах сборки качаем ОДИН раз (кэш): повторные быстрые
    запросы того же файла отбивались источником и слайд оставался без фото."""
    if not src:
        return None
    hit = _FETCH_CACHE.get(src)
    if hit is not None and (hit is False or Path(hit).is_file()):
        return None if hit is False else hit
    try:
        if src.startswith("file://"):
            p = Path(urllib.parse.unquote(src[7:]))
            return _to_supported(p, workdir, idx) if p.is_file() else None
        if src.startswith("data:"):
            head, _, b64 = src.partition(",")
            ext = ".png" if "png" in head else (".jpg" if "jpe" in head or "jpg" in head else ".png")
            dst = workdir / f"cv_{idx}{ext}"
            dst.write_bytes(base64.b64decode(b64))
        else:
            ext = _ext_of(src)
            dst = workdir / f"cv_{idx}{ext}"
            req = urllib.request.Request(src, headers={"User-Agent": _UA})
            with urllib.request.urlopen(req, timeout=30) as r:
                dst.write_bytes(r.read())
        out = _to_supported(dst, workdir, idx)
        _FETCH_CACHE[src] = out if out else False
        return out
    except Exception:
        _FETCH_CACHE[src] = False
        return None


# ──────────────────────────────────────────────────────────────────────────
#  Примитивы вёрстки
# ──────────────────────────────────────────────────────────────────────────
def _rect(slide, x, y, w, h, fill=None, alpha=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
        if alpha is not None:
            _set_alpha(sp.fill.fore_color, alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line)
        sp.line.width = Pt(line_w)
    return sp


def _set_alpha(fore, pct_opaque):
    """Непрозрачность заливки в процентах (100 = полностью непрозрачно)."""
    srgb = fore._xFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    el = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct_opaque * 1000))})
    srgb.append(el)


def _txbox(slide, x, y, w, h, anchor="t"):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    return tf


def _set_spc(run, em, font_px):
    """Разрядка (letter-spacing) через XML-атрибут spc (в сотых pt)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(round(em * font_px * 50))))


def _para(tf, first, text, size_px, color, font=GOLOS, bold=False,
          before_px=0, spc_em=None, line=None, align=None, upper=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if before_px:
        p.space_before = PT(before_px)
    if line:
        p.line_spacing = line
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = (text or "").upper() if upper else (text or "")
    r.font.name = font
    r.font.size = PT(size_px)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    if spc_em:
        _set_spc(r, spc_em, size_px)
    return p, r


def _img_contain(slide, path, x, y, w, h):
    """Вписать фото в область (contain), по центру, с сохранением пропорций.
    Возвращает фактический прямоугольник ОТРИСОВАННОГО фото в px
    (draw_x, draw_y, draw_w, draw_h) — по нему выравниваются подписи; None при
    отсутствии/ошибке фото."""
    if not path:
        return None
    try:
        pic = slide.shapes.add_picture(str(path), E(x), E(y))
        nw, nh = pic.width, pic.height
        tw, th = E(w), E(h)
        scale = min(tw / nw, th / nh)
        fw, fh = int(nw * scale), int(nh * scale)
        pic.width, pic.height = fw, fh
        left = int(E(x) + (tw - fw) / 2)
        # нижняя грань фото — на общую линию ячеек (правка Вашика 30.07:
        # несколько фото на слайде были на разных уровнях)
        top = int(E(y) + (th - fh))
        pic.left, pic.top = Emu(left), Emu(top)
        return (left / PXE, top / PXE, fw / PXE, fh / PXE)
    except Exception:
        return None


def _caption_bounds(cell_x, cell_w, drawn):
    """Границы подписи по ОТРИСОВАННОМУ фото: левый край = левому краю фото,
    ширина = ширине фото. Узкое фото → минимум 60% ширины ячейки, блок подписи
    центрируется по фото (в пределах ячейки). Нет фото → вся ширина ячейки."""
    if not drawn:
        return cell_x, cell_w
    dx, _, dw, _ = drawn
    min_w = 0.6 * cell_w
    if dw >= min_w:
        return dx, dw
    x = (dx + dw / 2) - min_w / 2
    x = max(cell_x, min(x, cell_x + cell_w - min_w))
    return x, min_w


def _img_fill_h(slide, path, cell_x, cell_w, y, h):
    """Фото РОВНО высотой h (широкие кадры кропятся по бокам — object-fit:cover
    по горизонтали), по центру ячейки. Возвращает прямоугольник как _img_contain.
    Нужен duo-clean: фото на слайде должны быть строго одного размера по высоте."""
    if not path:
        return None
    try:
        pic = slide.shapes.add_picture(str(path), E(cell_x), E(y))
        a = pic.width / float(pic.height)          # природный аспект
        if h * a <= cell_w:                        # влезает по ширине — полная высота
            w, ph, py = h * a, h, y
        else:
            # шире ячейки — вписываем целиком БЕЗ кропа («фото обрезаны», 27.07):
            # ширина = ячейка, высота меньше, центрируем в боксе
            w = cell_w
            ph = cell_w / a
            py = y + (h - ph) / 2.0
        left = cell_x + (cell_w - w) / 2.0
        pic.left, pic.top = E(left), E(py)
        pic.width, pic.height = E(w), E(ph)
        return (left, py, w, ph)
    except Exception:
        return None


def _img_cover(slide, path, x, y, w, h):
    """Фон object-fit:cover — кроп по осям, чтобы заполнить область без искажений."""
    if not path:
        return
    try:
        pic = slide.shapes.add_picture(str(path), E(x), E(y))
        nw, nh = pic.width, pic.height
        a = nw / float(nh)
        t = float(w) / float(h)
        if a > t:                       # шире области — режем бока
            crop = (1 - t / a) / 2.0
            pic.crop_left = pic.crop_right = crop
        else:                           # выше области — режем верх/низ
            crop = (1 - a / t) / 2.0
            pic.crop_top = pic.crop_bottom = crop
        pic.left, pic.top = E(x), E(y)
        pic.width, pic.height = E(w), E(h)
    except Exception:
        pass


def _logo(slide, path, x, y, h, right=None):
    """Вставить логотип заданной высоты (px), сохранив пропорции. Если задан
    right — прижать правым краём к этой X-координате. True, если файл вставлен;
    False (нет файла/ошибка) — вызывающий рисует прежний текст «STARGIFT»."""
    try:
        if not Path(path).exists():
            return False
        pic = slide.shapes.add_picture(str(path), E(x), E(y), height=E(h))
        if right is not None:
            pic.left = Emu(int(E(right) - pic.width))
        return True
    except Exception:
        return False


def _bg(slide, hexc):
    _rect(slide, 0, 0, 1920, 1080, fill=hexc)


def _header(slide, kicker, title, kick_color, title_color,
            x=110, w=None, title_size=70, kick_upper=False):
    """Кикер + заголовок вверху слайда. Возвращает y низа шапки (старт контента)."""
    if w is None:
        w = 1920 - 2 * x
    tf = _txbox(slide, x, 90, w, 44)
    _para(tf, True, kicker, 26, kick_color, spc_em=0.26, upper=kick_upper)
    tf2 = _txbox(slide, x, 138, w, title_size * 1.6)
    _para(tf2, True, title, title_size, title_color, font=PRATA)
    return 138 + title_size + 82   # ≈ отступ margin-bottom 50 из эталона


# ──────────────────────────────────────────────────────────────────────────
#  Двухуровневая подпись карточки + цена в фиксированной нижней позиции
# ──────────────────────────────────────────────────────────────────────────
def _caption_parts(lot):
    """(имя_подписанта, краткое_описание, полное_имя).
    Если в лоте ЕСТЬ ключ person (даже пустой) — resolveSlides явно управляет
    уровнями, доверяем полям как есть: person=="" означает «имя не рендерить»
    (верхним уровнем станет subtitle). Сплит name по « — » применяем ТОЛЬКО
    когда ключа person нет вовсе (старые данные)."""
    subtitle = _clean(lot.get("subtitle"))
    name = _clean(lot.get("name"))
    if "person" in lot:
        return _clean(lot.get("person")), subtitle, name
    if " — " in name:
        a, b = name.split(" — ", 1)
        b = b.strip()
        if b:
            b = b[0].upper() + b[1:]
        return a.strip(), b, name
    return name, "", name


def _cap_lines(txt, size_px, width_px, maxl):
    """Оценка числа строк текста в боксе (для расчёта высот и обрезки desc)."""
    if not txt:
        return 0
    import math
    cpl = max(1, int(width_px / (size_px * 0.52)))
    return min(maxl, max(1, math.ceil(len(txt) / cpl)))


def _zone_h(name_size, sub_size, price_size, with_desc, desc_size=0):
    """Высота нижней текстовой зоны под три уровня (имя 2 стр + подзаголовок 2 стр
    [+ desc 3 стр] + строка цены). Худший случай — чтобы не было наездов."""
    h = name_size * 1.26 * 2 + 8            # имя до 2 строк
    h += sub_size * 1.3 * 2 + 8             # подзаголовок до 2 строк
    if with_desc:
        h += 12 + desc_size * 1.42 * 3      # desc до 3 строк
    h += 14 + price_size * 1.5             # строка цены
    return int(h)


def _render_caption(slide, x, w, region_top, price_y, lot, *,
                    name_size, sub_size, desc_size, price_size,
                    name_c, sub_c, desc_c, price_c, cert_c,
                    price_align=PP_ALIGN.LEFT, with_desc=True, upper_cert=True):
    """Имя (крупно) + подзаголовок (мельче, приглушённо) + desc (ещё мельче,
    обрезка многоточием под доступную высоту). Цена — отдельным боксом на
    фиксированной Y (price_y), поэтому НЕ зависит от объёма текста выше.
    Имя и подзаголовок не режем (перенос по словам), режем только desc."""
    person, subtitle, _ = _caption_parts(lot)
    # если имени нет (person==""), крупным уровнем становится подзаголовок
    if person:
        top_txt, second_txt = person, subtitle
    else:
        top_txt, second_txt = subtitle, ""
    zone = max(10, price_y - 8 - region_top)
    tf = _txbox(slide, x, region_top, w, zone, "t")
    first = True
    if top_txt:
        _para(tf, first, top_txt, name_size, name_c, bold=True, line=1.26); first = False
    if second_txt:
        _para(tf, first, second_txt, sub_size, sub_c,
              before_px=(0 if first else 8), line=1.3); first = False
    desc = _clean(lot.get("desc"))
    if with_desc and desc:
        nl = _cap_lines(top_txt, name_size, w, 2)
        sl = _cap_lines(second_txt, sub_size, w, 2)
        used = nl * name_size * 1.26 + (8 + sl * sub_size * 1.3 if second_txt else 0)
        avail = zone - used - 12
        dlh = desc_size * 1.42
        dmax = int(avail / dlh) if dlh else 0
        if dmax >= 1:
            cpl = max(1, int(w / (desc_size * 0.52)))
            budget = dmax * cpl
            if len(desc) > budget:
                desc = _fit_sent(desc, budget)
            _para(tf, first, desc, desc_size, desc_c,
                  before_px=(0 if first else 12), line=1.42); first = False
    # цена — фиксированная нижняя позиция ячейки
    pt = _txbox(slide, x, price_y, w, price_size * 1.6, "t")
    _para(pt, True, _fmt_rub(lot["price"]), price_size, price_c, font=PRATA, align=price_align)
    # инфо-блок: спецификация и наличие двумя строками, без «·» (правка 28.07)
    ca = PP_ALIGN.RIGHT if price_align == PP_ALIGN.LEFT else PP_ALIGN.LEFT
    fs = max(18, int(price_size * 0.6))
    ct = _txbox(slide, x, price_y, w, fs * 3.4, "t")
    first = True
    if _spec_line(lot):
        _para(ct, True, _spec_line(lot), fs, cert_c, spc_em=0.02, align=ca, upper=False)
        first = False
    _para(ct, first, _avail_line(lot), fs, cert_c, spc_em=0.02, align=ca,
          upper=False, before_px=(0 if first else 4))


def _photo_cell(slide, x, w, cell_top, cell_bottom, lot, imgs, *, style,
                tall=False, pad_top=0, pad_bottom=0):
    """Фото сверху (contain, фикс-бокс) + двухуровневая подпись снизу + цена на
    фикс-Y. Возвращает price_y — общий для всех ячеек слайда (одинаковые
    cell_top/cell_bottom → одинаковый price_y → цены в одну горизонталь)."""
    with_desc = not tall
    zh = _zone_h(style["name"], style["sub"], style["price"], with_desc, style["desc"])
    region_bottom = cell_bottom - pad_bottom
    region_top = region_bottom - zh
    photo_top = cell_top + pad_top
    photo_bottom = region_top - 16
    drawn = _img_contain(slide, imgs.get(id(lot)), x, photo_top, w, max(20, photo_bottom - photo_top))
    cap_x, cap_w = _caption_bounds(x, w, drawn)     # подпись по краям отрисованного фото
    price_y = region_bottom - int(style["price"] * 1.5)
    _render_caption(slide, cap_x, cap_w, region_top, price_y, lot,
                    name_size=style["name"], sub_size=style["sub"],
                    desc_size=style["desc"], price_size=style["price"],
                    name_c=style["name_c"], sub_c=style["sub_c"],
                    desc_c=style["desc_c"], price_c=style["price_c"],
                    cert_c=style["cert_c"], price_align=style.get("price_align", PP_ALIGN.LEFT),
                    with_desc=with_desc, upper_cert=style.get("upper_cert", True))
    return price_y


# ──────────────────────────────────────────────────────────────────────────
#  Слайды
# ──────────────────────────────────────────────────────────────────────────
def _cap_center(slide, x, w, y, lot, heading, max_h=340):
    """Центрированная подпись под фото (эталон «Скала» 30.07): имя — контекст,
    desc, цена золотом, спецификация и наличие мелким капсом. Все строки по центру."""
    person, subtitle, _ = _caption_parts(lot)
    name = person or subtitle or ""
    sub = subtitle if person else ""
    tf = _txbox(slide, x, y, w, 208, "t")
    _para(tf, True, name, 26, TXT_D, font=PRATA, line=1.35, align=PP_ALIGN.CENTER)
    if sub:
        # имя подписанта и название экспоната НИКОГДА не в одну строку (правка 30.07)
        _para(tf, False, sub, 20, "6A6A6A", before_px=6, line=1.4, align=PP_ALIGN.CENTER)
    desc = _clean(lot.get("desc"))
    if desc:
        _para(tf, False, _fit_sent(desc, 170), 19, "6A6A6A", before_px=10, line=1.45,
              align=PP_ALIGN.CENTER)
    # цена/спецификация/наличие — ОТДЕЛЬНЫМ боксом на фиксированном уровне
    # (правка Вашика 30.07: строки соседних карточек совпадают по горизонтали)
    pb = _txbox(slide, x, y + 214, w, 126, "t")
    _para(pb, True, _fmt_rub(lot.get("price", 0)), 26, GOLD_L, font=PRATA,
          align=PP_ALIGN.CENTER)
    if _spec_line(lot):
        _para(pb, False, _spec_line(lot), 16, "9A9A9A", before_px=10, spc_em=0.14,
              upper=True, align=PP_ALIGN.CENTER)
    _para(pb, False, _avail_line(lot), 16, "9A9A9A", before_px=4, spc_em=0.14,
          upper=True, align=PP_ALIGN.CENTER)


def _sentence_case(t):
    t = (t or "").strip()
    if not t:
        return ""
    import re as _re
    return t if _re.search(r"[a-zа-яё]", t) else (t[:1].upper() + t[1:].lower())


def _light_head(slide, kicker, title, count=0):
    """Шапка светлого слайда: заголовок + серый подзаголовок-эпитет, счётчик справа."""
    ht = _txbox(slide, 110, 84, 1300, 130)
    _para(ht, True, title or "", 46, TXT_D, font=PRATA, line=1.1)
    if kicker:
        _para(ht, False, _sentence_case(kicker), 22, MUTED, before_px=10, line=1.35)
    if count:
        kt = _txbox(slide, 1100, 100, 710, 40)
        _para(kt, True, "%d %s" % (count, _plural(count, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ")),
              18, "9A9A9A", spc_em=0.22, upper=True, align=PP_ALIGN.RIGHT)


def _slide_cover(slide, p, imgs):
    _bg(slide, DARK)
    _img_cover(slide, imgs.get("cover"), 0, 0, 1920, 1080)
    _rect(slide, 0, 0, 1920, 1080, fill=DARK, alpha=52)      # фолбэк-затемнение
    # левый градиент плотнее — размывает границу фото и чёрного поля (правка 27.07)
    _scrim(slide, 0, 0, 1920, 1080, ang_deg=0,
           stops=[(0, 90), (38000, 52), (74000, 0), (100000, 0)])
    # верхняя строка: слева логотип (светлый), справа кикер
    if not _logo(slide, LOGO_WHITE, 110, 44, 230):
        lt = _txbox(slide, 110, 90, 900, 44)
        _para(lt, True, "STARGIFT", 27, "D0D0D0", spc_em=0.28)
    rt = _txbox(slide, 910, 90, 900, 44, "t")
    _para(rt, True, p.get("kicker", ""), 27, "D0D0D0", spc_em=0.28, align=PP_ALIGN.RIGHT)
    # нижний блок (заголовок / подзаголовок / счётчик), прижат к низу
    tf = _txbox(slide, 110, 380, 1300, 620, "b")
    _para(tf, True, p.get("title", ""), 118, TXT_L, font=PRATA, line=1.04)
    if p.get("sub"):
        _para(tf, False, p["sub"], 35, TXT_L2, before_px=36, line=1.45)
    # строку-счётчик показываем только если count реально задан (>0),
    # иначе «0 ЭКСПОНАТОВ» — пустой мусор
    cnt = int(p.get("count", 0) or 0)
    if cnt > 0:
        tail = "%d %s" % (cnt, _plural(cnt, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ"))
        if p.get("certs"):
            tail += ", " + p["certs"]
        _para(tf, False, tail, 27, GOLD, before_px=36, spc_em=0.24)
    elif p.get("certs"):
        _para(tf, False, p["certs"], 27, GOLD, before_px=36, spc_em=0.24)
    # подпись источника архивного кадра обложки — мелко справа внизу
    if p.get("credit"):
        cr = _txbox(slide, 910, 1010, 900, 40, "b")
        _para(cr, True, p["credit"].upper(), 18, "8D7A58",
              spc_em=0.14, align=PP_ALIGN.RIGHT)


def _slide_divider(slide, p, imgs):
    _bg(slide, DARK)
    _img_cover(slide, imgs.get("bg"), 0, 0, 1920, 1080)
    # затемнение слева (текстовая зона) сильнее — два прямоугольника-фолбэка
    _rect(slide, 0, 0, 1920, 1080, fill=DARK, alpha=38)
    _rect(slide, 0, 0, 1150, 1080, fill=DARK, alpha=52)
    tf = _txbox(slide, 110, 0, 1050, 1080, "m")
    _para(tf, True, "РАЗДЕЛ " + str(p.get("roman", "")), 27, GOLD, spc_em=0.3)
    _para(tf, False, p.get("title", ""), 92, TXT_L, font=PRATA, before_px=40, line=1.06)
    if p.get("subtitle"):
        _para(tf, False, p["subtitle"], 42, "E6E6E6", before_px=40, line=1.5)
    cnt = int(p.get("count", 0) or 0)
    tail = "%d %s, ОТ %s" % (cnt, _plural(cnt, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ"),
                              _fmt_rub(p.get("min", 0)))
    _para(tf, False, tail, 27, GOLD, before_px=40, spc_em=0.24)
    if p.get("credit"):
        cr = _txbox(slide, 910, 1010, 900, 40, "b")
        _para(cr, True, p["credit"].upper(), 18, "8D7A58", spc_em=0.14, align=PP_ALIGN.RIGHT)


def _slide_contents(slide, p, imgs):
    _bg(slide, PAPER)
    cards = p.get("cards", [])
    total = int(p.get("total", 0) or 0)
    cnt = int(p.get("count", 0) or 0)
    # верхняя строка
    top = _txbox(slide, 110, 96, 1700, 40)
    _para(top, True, "КАТАЛОГ, %d %s" % (cnt, _plural(cnt, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ")),
          26, MUTED, spc_em=0.26)
    # правый верх — тёмный логотип на светлом фоне (фолбэк — текст)
    if not _logo(slide, LOGO_DARK, 0, 74, 72, right=1810):
        tr = _txbox(slide, 110, 96, 1700, 40)
        _para(tr, True, "STARGIFT", 26, MUTED, spc_em=0.26, align=PP_ALIGN.RIGHT)
    ht = _txbox(slide, 110, 140, 1700, 110)
    _para(ht, True, p.get("title") or "Состав собрания", 76, TXT_D, font=PRATA)
    area_y, area_b = 300, 940
    if len(cards) > 6:
        # компактный перечень в 2 колонки
        col_w, gap = (1700 - 80) / 2.0, 80
        half = (len(cards) + 1) // 2
        for i, s in enumerate(cards):
            col, row = (0, i) if i < half else (1, i - half)
            cy = area_y + row * 62
            cx = 110 + col * (col_w + gap)
            rowtf = _txbox(slide, cx, cy, col_w, 56)
            _para(rowtf, True, "%s. %s" % (s.get("roman", ""), s.get("title", "")), 26, TXT_D)
            prtf = _txbox(slide, cx, cy, col_w, 56)
            _para(prtf, True, "%d, от %s" % (int(s.get("count", 0)), _fmt_rub(s.get("min", 0))),
                  26, GOLD_L, bold=True, align=PP_ALIGN.RIGHT)
    else:
        cols = min(3, len(cards)) or 1
        gap = 40
        cw = (1700 - gap * (cols - 1)) / cols
        import math
        rows = max(1, math.ceil(len(cards) / cols))
        ch = (area_b - area_y - gap * (rows - 1)) / rows
        for i, s in enumerate(cards):
            col, row = i % cols, i // cols
            cx = 110 + col * (cw + gap)
            cy = area_y + row * (ch + gap)
            _rect(slide, cx, cy, cw, ch, fill=CARD_L, line=CARD_L_BD, line_w=1.0)
            tf = _txbox(slide, cx + 42, cy + 40, cw - 84, ch - 80)
            _para(tf, True, "%s. %s" % (s.get("roman", ""), (s.get("title", "")).upper()),
                  25, MUTED, spc_em=0.22)
            _para(tf, False, s.get("title", ""), 36, TXT_D, font=PRATA, before_px=22, line=1.2)
            bt = _txbox(slide, cx + 42, cy + ch - 80, cw - 84, 60, "b")
            n = int(s.get("count", 0))
            _para(bt, True, "%d %s, от %s" % (n, _plural(n, "экспонат", "экспоната", "экспонатов"),
                  _fmt_rub(s.get("min", 0))), 29, GOLD_L, bold=True)
    # итог
    ft = _txbox(slide, 110, 960, 1700, 60)
    _para(ft, True, "ИТОГ СОБРАНИЯ", 27, MUTED, spc_em=0.24)
    fr = _txbox(slide, 110, 960, 1700, 60)
    _para(fr, True, _fmt_rub(total), 27, TXT_D, font=PRATA, align=PP_ALIGN.RIGHT)
    _rect(slide, 110, 952, 1700, 2, fill=LINE_L)


def _slide_premium_solo(slide, p, imgs):
    # лот приходит в массиве lots (по спеке), с фолбэком на params.lot
    lots = p.get("lots") or ([p["lot"]] if p.get("lot") else [])
    lot = lots[0] if lots else {}
    light = bool(p.get("light"))
    bg = PAPER if light else DARK
    txt_c = TXT_D if light else TXT_L
    sub_c = DESC_L if light else DESC_D
    ds_c = "3A3A3A" if light else TXT_L2
    pr_c = GOLD_L if light else GOLD_HL
    kk_c = MUTED if light else GOLD
    _bg(slide, bg)
    _img_contain(slide, imgs.get(id(lot)), 60, 70, 763, 940)
    person, subtitle, _ = _caption_parts(lot)
    # имени нет → крупным уровнем (Prata) становится подзаголовок
    top_txt, second_txt = (person, subtitle) if person else (subtitle, "")
    desc = _clean(lot.get("desc"))
    tf = _txbox(slide, 973, 90, 837, 900, "m")
    _para(tf, True, p.get("kicker", "ПРЕМИУМ-ЛОТ РАЗДЕЛА"), 26, kk_c, spc_em=0.26)
    # два уровня: имя подписанта крупно (Prata) + краткое описание мельче
    if top_txt:
        _para(tf, False, top_txt, 60, txt_c, font=PRATA, before_px=34, line=1.08)
    if second_txt:
        _para(tf, False, second_txt, 33, sub_c, before_px=20, line=1.28)
    if desc:
        _para(tf, False, desc, 31, ds_c, before_px=28, line=1.5)
    _para(tf, False, _fmt_rub(lot.get("price", 0)), 33, pr_c, font=PRATA, before_px=32)
    if lot.get("size"):
        _para(tf, False, lot["size"], 27, sub_c, before_px=26)
    cert = (lot.get("cert") or "").strip()
    if cert and cert != "Provenance":
        _para(tf, False, cert, 27, sub_c, before_px=14)
    _para(tf, False, _avail_line(lot), 27, sub_c, before_px=14)


_STYLE_GRID4 = {"name": 27, "sub": 21, "desc": 21, "price": 25,
                "name_c": TXT_D, "sub_c": MUTED, "desc_c": DESC_L,
                "price_c": GOLD_L, "cert_c": MUTED}


def _slide_grid4(slide, p, imgs):
    """Крест 2×2 · вечерний (q01 ✓): тёмный фон, 4 экспоната сеткой 2×2, крупные фото,
    имя+тип слева и цена (золото) справа по краям фото. desc/cert не выводятся."""
    light = bool(p.get("light"))
    _bg(slide, PAPER if light else DARK)
    g_txt = TXT_D if light else TXT_L
    g_sub = DESC_L if light else DESC_D
    g_pr = GOLD_L if light else GOLD_HL
    lots = (p.get("lots") or [])[:4]
    if light:
        # светлая сетка — подпись по центру (эталон «Скала»)
        _light_head(slide, p.get("kicker", ""), p.get("heading") or p.get("title") or "Экспонаты раздела", len(lots))
        n = len(lots) or 1
        gap = 44
        cw = (1740 - gap * (n - 1)) / n
        for i, lot in enumerate(lots):
            cx = 90 + i * (cw + gap)
            _img_contain(slide, imgs.get(id(lot)), cx, 280, cw, 460)
            _cap_center(slide, cx, cw, 770, lot, p.get("heading") or "")
        return
    # кикер-плашка убрана (правка Вашика к q01)
    # сжатые поля/шапка — фото максимально крупные (правка 27.07 «очень маленькие фото»)
    ht = _txbox(slide, 90, 66, 1740, 84)
    _para(ht, True, (p.get("heading") or p.get("title") or "Экспонаты раздела"), 62, g_txt, font=PRATA)
    n = len(lots) or 1
    cols = n                        # все в один ряд — фото вдвое крупнее (правка 27.07)
    import math
    rows = 1
    top, bottom = 196, 1020
    gx, gy = 44, 26
    cw = (1740 - gx * (cols - 1)) / cols
    ch = (bottom - top - gy * (rows - 1)) / rows
    photo_h = ch - 118
    price_w = 160          # узкие колонки 1×4 — имени нужна вся ширина ячейки
    for i, lot in enumerate(lots):
        col, row = i % cols, i // cols
        cx = 90 + col * (cw + gx)
        cy = top + row * (ch + gy)
        _img_fill_h(slide, imgs.get(id(lot)), cx, cw, cy, photo_h)
        cap_x, cap_w = cx, cw          # подпись по ширине ячейки, не отрисованного фото
        ty = cy + photo_h + 18
        person, subtitle, _ = _caption_parts(lot)
        top_txt, second_txt = (person, subtitle) if person else (subtitle, "")
        nb = _txbox(slide, cap_x, ty, max(120, cap_w - price_w), 90, "t")
        first = True
        if top_txt:
            _para(nb, True, top_txt, 26, g_txt, bold=True, line=1.2); first = False
        if second_txt:
            _para(nb, first, second_txt, 21, g_sub, before_px=(0 if first else 5), line=1.2)
        # цена на нижней линии СВОЕГО ряда (унификация 27.07)
        row_bottom = cy + ch - 6
        pr = _txbox(slide, cap_x + cap_w - price_w, row_bottom - 44, price_w, 44, "b")
        _para(pr, True, _fmt_rub(lot["price"]), 26, g_pr, font=PRATA, align=PP_ALIGN.RIGHT)


def _slide_pair_dark(slide, p, imgs):
    _bg(slide, DARK)
    lots = (p.get("lots") or [])[:2]
    top = _header(slide, p.get("kicker", ""), (p.get("heading") or p.get("title") or "Экспонаты раздела"),
                  GOLD, TXT_L)
    gap = 44
    cw = (1700 - gap) / 2.0
    ch = 1000 - top
    pad = 32
    for i, lot in enumerate(lots):
        cx = 110 + i * (cw + gap)
        _rect(slide, cx, top, cw, ch, fill=CARD_D)
        # фото фикс. 392, ниже — двухуровневая подпись; цена на фикс-Y снизу карточки
        drawn = _img_contain(slide, imgs.get(id(lot)), cx + pad, top + pad, cw - 2 * pad, 440)
        cap_x, cap_w = _caption_bounds(cx + pad, cw - 2 * pad, drawn)
        region_top = top + pad + 392 + 22
        price_y = top + ch - pad - int(26 * 1.5)
        _render_caption(slide, cap_x, cap_w, region_top, price_y, lot,
                        name_size=30, sub_size=24, desc_size=24, price_size=26,
                        name_c=TXT_L, sub_c=DESC_D, desc_c=DESC_D, price_c=GOLD_HL,
                        cert_c=MUTED, price_align=PP_ALIGN.LEFT)


_STYLE_PAIR_LIGHT = {"name": 31, "sub": 24, "desc": 24, "price": 26,
                     "name_c": TXT_D, "sub_c": MUTED, "desc_c": DESC_L,
                     "price_c": GOLD_L, "cert_c": MUTED, "price_align": PP_ALIGN.RIGHT}


def _slide_pair_light(slide, p, imgs):
    _bg(slide, PAPER)
    lots = (p.get("lots") or [])[:2]
    _light_head(slide, p.get("kicker", ""), p.get("heading") or p.get("title") or "Экспонаты раздела", len(lots))
    tall = _all_desc_empty(lots)
    gap = 80
    cw = (1700 - gap) / 2.0
    ph = 430 if not tall else 520
    for i, lot in enumerate(lots):
        cx = 110 + i * (cw + gap)
        _img_contain(slide, imgs.get(id(lot)), cx, 290, cw, ph - 50)
        _cap_center(slide, cx, cw, 290 + (ph - 50) + 30, lot, p.get("heading") or "")


def _slide_cards_side(slide, p, imgs):
    _bg(slide, PAPER)
    lots = (p.get("lots") or [])[:2]
    top = _header(slide, p.get("kicker", ""), (p.get("heading") or p.get("title") or "Экспонаты раздела"),
                  MUTED, TXT_D)
    gap = 44
    cw = (1700 - gap) / 2.0
    ch = 1000 - top
    for i, lot in enumerate(lots):
        cx = 110 + i * (cw + gap)
        _rect(slide, cx, top, cw, ch, fill=CARD_L, line=CARD_L_BD, line_w=1.0)
        pad = 44
        iw = (cw - 2 * pad - 40) * 0.42
        _img_contain(slide, imgs.get(id(lot)), cx + pad, top + pad, iw, ch - 2 * pad)
        tx = cx + pad + iw + 40
        tw = cx + cw - pad - tx
        # фото слева, справа — двухуровневая подпись; цена на фикс-Y внизу карточки
        region_top = top + pad
        price_y = top + ch - pad - int(27 * 1.5)
        _render_caption(slide, tx, tw, region_top, price_y, lot,
                        name_size=31, sub_size=25, desc_size=27, price_size=27,
                        name_c=TXT_D, sub_c=MUTED, desc_c=DESC_L, price_c=GOLD_L,
                        cert_c=MUTED, price_align=PP_ALIGN.LEFT)


def _slide_trio(slide, p, imgs, dark):
    bg = DARK if dark else PAPER
    text_c = TXT_L if dark else TXT_D
    card_bg = CARD_D if dark else CARD_L
    card_bd = None if dark else CARD_L_BD
    kick_c = GOLD if dark else MUTED
    desc_c = DESC_D if dark else DESC_L
    price_c = GOLD_HL if dark else GOLD_L
    _bg(slide, bg)
    lots = (p.get("lots") or [])[:3]
    tall = _all_desc_empty(lots)
    if not dark:
        # светлое трио — чистый белый, подпись по центру (эталон «Скала»)
        _light_head(slide, p.get("kicker", ""), p.get("heading") or p.get("title") or "Экспонаты раздела", len(lots))
        n = len(lots) or 1
        gap = 64
        cw = (1700 - gap * (n - 1)) / n
        ph = 380 if not tall else 470
        for i, lot in enumerate(lots):
            cx = 110 + i * (cw + gap)
            _img_contain(slide, imgs.get(id(lot)), cx, 300, cw, ph - 60)
            _cap_center(slide, cx, cw, 300 + (ph - 60) + 30, lot, p.get("heading") or "")
        return
    kt = _txbox(slide, 90, 90, 1740, 44)
    _para(kt, True, p.get("kicker", ""), 26, kick_c, spc_em=0.26)
    ht = _txbox(slide, 90, 138, 1740, 100)
    _para(ht, True, (p.get("heading") or p.get("title") or "Экспонаты раздела"), 70, text_c, font=PRATA)
    top, bottom = 270, 1000
    n = len(lots) or 1
    gap = 44
    cw = (1740 - gap * (n - 1)) / n
    ch = bottom - top
    padx, padt, padb = 38, 40, 42
    with_desc = not tall
    # зона текста снизу карточки (имя+подзаголовок[+desc]+цена); фото — верх
    zone = _zone_h(31, 23, 30, with_desc, 23)
    region_bottom = top + ch - padb
    region_top = region_bottom - zone
    price_y = region_bottom - int(30 * 1.5)     # фикс-Y цены, общий для всех карточек
    for i, lot in enumerate(lots):
        cx = 90 + i * (cw + gap)
        _rect(slide, cx, top, cw, ch, fill=card_bg, line=card_bd, line_w=1.0)
        drawn = _img_contain(slide, imgs.get(id(lot)), cx + padx, top + padt,
                             cw - 2 * padx, (region_top - 16) - (top + padt))
        cap_x, cap_w = _caption_bounds(cx + padx, cw - 2 * padx, drawn)
        _render_caption(slide, cap_x, cap_w, region_top, price_y, lot,
                        name_size=31, sub_size=23, desc_size=23, price_size=30,
                        name_c=text_c, sub_c=(MUTED if not dark else DESC_D),
                        desc_c=desc_c, price_c=price_c, cert_c=MUTED,
                        price_align=PP_ALIGN.LEFT, with_desc=with_desc)


def _slide_duo_clean(slide, p, imgs, dark):
    """Два экспоната без описаний — витрина «duo-clean»: крупные фото прямо на
    фоне (БЕЗ карточек и рамок, БЕЗ теней), под каждым — одна строка подписи по
    краям фото: слева имя (27) + тип-подзаголовок (22 muted), справа цена (Prata
    30) на уровне первой строки имени. desc и cert не выводятся никогда."""
    bg = DARK if dark else PAPER
    text_c = TXT_L if dark else TXT_D
    kick_c = GOLD if dark else MUTED
    sub_c = DESC_D if dark else MUTED
    price_c = GOLD_HL if dark else GOLD_L
    _bg(slide, bg)
    lots = (p.get("lots") or [])[:2]
    if not dark:
        _light_head(slide, p.get("kicker", ""), p.get("heading") or p.get("title") or "Экспонаты раздела", len(lots))
    else:
        kt = _txbox(slide, 90, 90, 1740, 44)
        _para(kt, True, p.get("kicker", ""), 26, kick_c, spc_em=0.26)
        ht = _txbox(slide, 90, 138, 1740, 100)
        _para(ht, True, (p.get("heading") or p.get("title") or "Экспонаты раздела"),
              70, text_c, font=PRATA)
    light = not dark
    top = 270 if dark else 300
    n = len(lots) or 1
    gap = 60 if dark else 80
    cw = (1740 - gap * (n - 1)) / n
    photo_h = 560 if dark else 460
    price_w = 250
    for i, lot in enumerate(lots):
        cx = 90 + i * (cw + gap)
        drawn = _img_fill_h(slide, imgs.get(id(lot)), cx, cw, top, photo_h)
        cap_x, cap_w = _caption_bounds(cx, cw, drawn)
        photo_bottom = (drawn[1] + drawn[3]) if drawn else (top + photo_h)
        cy = photo_bottom + (24 if dark else 30)
        if light:
            _cap_center(slide, cx, cw, cy, lot, p.get("heading") or "")
            continue
        person, subtitle, _ = _caption_parts(lot)
        top_txt, second_txt = (person, subtitle) if person else (subtitle, "")
        nb = _txbox(slide, cap_x, cy, max(120, cap_w - price_w), 120, "t")
        first = True
        if top_txt:
            _para(nb, True, top_txt, 27, text_c, bold=True, line=1.22)
            first = False
        if second_txt:
            _para(nb, first, second_txt, 22, sub_c, before_px=(0 if first else 6), line=1.25)
            first = False
        if _spec_line(lot):
            _para(nb, first, _spec_line(lot), 19, sub_c, before_px=(0 if first else 6), line=1.25)
            first = False
        _para(nb, first, _avail_line(lot), 19, sub_c, before_px=(0 if first else 4), line=1.25)
        # цена на ОБЩЕЙ нижней линии слайда (унификация 27.07)
        pb = _txbox(slide, cap_x + cap_w - price_w, 952, price_w, 48, "b")
        _para(pb, True, _fmt_rub(lot.get("price", 0)), 30, price_c, font=PRATA,
              align=PP_ALIGN.RIGHT)


def _slide_instock(slide, p, imgs):
    """«В наличии» в стиле duo-clean: 1–2 лота, фото КРУПНО прямо на тёмном фоне
    (строго одной высоты, _img_fill_h), под фото имя+тип слева и цена справа
    по краям фото. Строка «и ещё N…» — под заголовком, если есть extra."""
    light = bool(p.get("light"))
    _bg(slide, PAPER if light else DARK)
    i_txt = TXT_D if light else TXT_L
    i_sub = DESC_L if light else DESC_D
    i_pr = GOLD_L if light else GOLD_HL
    i_kk = MUTED if light else GOLD
    lots = (p.get("lots") or [])[:2]
    extra = int(p.get("extra", 0) or 0)
    kt = _txbox(slide, 90, 90, 1740, 44)
    _para(kt, True, p.get("kicker", "ГОТОВЫ К ПЕРЕДАЧЕ"), 26, i_kk, spc_em=0.26)
    ht = _txbox(slide, 90, 138, 1740, 100)
    _para(ht, True, p.get("title", "Экспонаты в наличии"), 70, i_txt, font=PRATA)
    top = 270
    if extra:
        et = _txbox(slide, 90, 254, 1740, 40)
        _para(et, True, "и ещё %d %s в наличии — в составе разделов" %
              (extra, _plural(extra, "экспонат", "экспоната", "экспонатов")), 26, i_sub, line=1.4)
        top = 324
    photo_h = (480 if extra else 520) if light else (540 if extra else 590)
    n = len(lots) or 1
    gap = 80 if light else 60
    cw = (1740 - gap * (n - 1)) / n
    price_w = 250
    for i, lot in enumerate(lots):
        cx = 90 + i * (cw + gap)
        drawn = _img_fill_h(slide, imgs.get(id(lot)), cx, cw, top, photo_h)
        cap_x, cap_w = _caption_bounds(cx, cw, drawn)
        photo_bottom = (drawn[1] + drawn[3]) if drawn else (top + photo_h)
        cy = photo_bottom + (28 if light else 24)
        if light:
            _cap_center(slide, cx, cw, cy, lot, p.get("title") or "")
            continue
        person, subtitle, _ = _caption_parts(lot)
        top_txt, second_txt = (person, subtitle) if person else (subtitle, "")
        nb = _txbox(slide, cap_x, cy, max(120, cap_w - price_w), 120, "t")
        first = True
        if top_txt:
            _para(nb, True, top_txt, 27, i_txt, bold=True, line=1.22)
            first = False
        if second_txt:
            _para(nb, first, second_txt, 22, i_sub,
                  before_px=(0 if first else 6), line=1.25)
            first = False
        if _spec_line(lot):
            _para(nb, first, _spec_line(lot), 19, i_sub, before_px=(0 if first else 6), line=1.25)
            first = False
        _para(nb, first, _avail_line(lot), 19, i_sub, before_px=(0 if first else 4), line=1.25)
        # цена на общей нижней линии слайда (унификация 27.07)
        pr = _txbox(slide, cap_x + cap_w - price_w, 952, price_w, 48, "b")
        _para(pr, True, _fmt_rub(lot["price"]), 30, i_pr, font=PRATA,
              align=PP_ALIGN.RIGHT)


def _slide_index(slide, p, imgs):
    _bg(slide, PAPER)
    lots = p.get("lots", [])
    total = int(p.get("total", 0) or 0)
    cnt = int(p.get("count", len(lots)) or len(lots))
    # кикер: если пришёл готовый текст cont — берём его; если cont=True — добавляем
    # пометку продолжения; иначе собираем из count как раньше
    cont = p.get("cont")
    if isinstance(cont, str) and cont.strip():
        kicker_line = cont.strip()
    else:
        kicker_line = "ПОЛНЫЙ ПЕРЕЧЕНЬ, %d %s" % (cnt, _plural(cnt, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ"))
        if cont:
            kicker_line += ", ПРОДОЛЖЕНИЕ"
    top = _txbox(slide, 110, 84, 1700, 40)
    _para(top, True, kicker_line, 25, MUTED, spc_em=0.26)
    tr = _txbox(slide, 110, 84, 1700, 40)
    _para(tr, True, "ЦЕНЫ В ₽", 25, MUTED, spc_em=0.26, align=PP_ALIGN.RIGHT)
    ht = _txbox(slide, 110, 128, 1700, 90)
    _para(ht, True, "Собрание «%s»" % (p.get("title", "")), 64, TXT_D, font=PRATA)

    # ── Перечень в два столбца ──
    # Надёжность прежде всего: каждая строка СТРОГО однострочная (эталон добивается
    # этого узкими колонками; здесь — обрезка имени + word_wrap=False), поэтому шаг
    # строки фиксированный и наездов не бывает ни при каких данных.
    import math
    area_y, area_b = 268, 936
    col_w, gap = (1700 - 80) / 2.0, 80
    name_w = col_w * 0.78                       # ширина под имя (остаток — под цену)
    row_h = 56                                  # фиксированный шаг строки
    max_rows = int((area_b - area_y) // row_h)  # сколько строк влезает в колонку (≈11)
    capacity = max_rows * 2                      # ёмкость двух колонок (≈22)

    all_lots = list(lots)
    overflow = len(all_lots) > capacity
    if overflow:
        # оставляем последнюю строку колонки под приписку «…и ещё N»
        rows_for_lots = max_rows - 1
        shown = all_lots[:rows_for_lots * 2]
        hidden = len(all_lots) - len(shown)
    else:
        shown = all_lots
        hidden = 0

    half = (len(shown) + 1) // 2                 # строк в первой колонке
    # запас по ширине имени: ~48 симв. на колонку при кегле 25px гарантируют одну строку
    NAME_MAX = 48
    for i, lot in enumerate(shown):
        col, row = (0, i) if i < half else (1, i - half)
        cx = 110 + col * (col_w + gap)
        cy = area_y + row * row_h
        note = "в наличии в Москве" if lot.get("instock") else "под заказ"
        nm = _clean(lot.get("name")) + "  (" + note + ")"
        if len(nm) > NAME_MAX:
            nm = nm[:NAME_MAX - 1].rstrip() + "…"
        lt = _txbox(slide, cx, cy, name_w, row_h)
        lt.word_wrap = False                    # никогда не переносить — только одна строка
        _para(lt, True, nm, 25, TXT_D, line=1.1)
        pt = _txbox(slide, cx + name_w, cy, col_w - name_w, row_h)
        _para(pt, True, _fmt_num(lot["price"]), 25, TXT_D, bold=True, align=PP_ALIGN.RIGHT)
        _rect(slide, cx, cy + row_h - 2, col_w, 1, fill=LINE_L2)
    if overflow:
        # приписка на свободной последней строке (правая колонка, снизу)
        ny = area_y + (max_rows - 1) * row_h
        nt = _txbox(slide, 110 + col_w + gap, ny, col_w, row_h)
        _para(nt, True, "…и ещё %d %s" %
              (hidden, _plural(hidden, "экспонат", "экспоната", "экспонатов")),
              25, MUTED, line=1.1)
    # итоговая строка — только если showTotal True (или флаг не задан → как раньше)
    if p.get("showTotal", True):
        _rect(slide, 110, 952, 1700, 3, fill=TXT_D)
        ft = _txbox(slide, 110, 962, 1700, 56)
        _para(ft, True, "%d %s, ВСЁ СОБРАНИЕ ЦЕЛИКОМ" %
              (cnt, _plural(cnt, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ")), 27, MUTED, spc_em=0.22)
        fr = _txbox(slide, 110, 962, 1700, 56)
        _para(fr, True, _fmt_rub(total), 26, TXT_D, font=PRATA, align=PP_ALIGN.RIGHT)


def _slide_quote(slide, p, imgs):
    _lots = p.get("lots") or ([p["lot"]] if p.get("lot") else [])
    lot = _lots[0] if _lots else {}
    _bg(slide, DARK)
    tf = _txbox(slide, 110, 100, 950, 880, "m")
    _para(tf, True, "СЛОВО ГЕРОЯ", 26, GOLD, spc_em=0.26)
    _para(tf, False, "«%s»" % p.get("quote", ""), 66, TXT_L, font=PRATA,
          before_px=44, line=1.18)
    if p.get("author"):
        _para(tf, False, "— " + p["author"], 29, DESC_D, before_px=44, spc_em=0.14)
    # правая колонка: фото + подпись
    _img_contain(slide, imgs.get(id(lot)), 1140, 130, 690, 640)
    bt = _txbox(slide, 1140, 800, 690, 200)
    _rect(slide, 1140, 790, 690, 1, fill=LINE_D)
    _para(bt, True, _clean(lot.get("name")), 28, TXT_L, bold=True, line=1.3, before_px=24)
    _para(bt, False, _fmt_rub(lot.get("price", 0)), 27, GOLD_HL, font=PRATA, before_px=12)
    ml = _meta_line(lot)
    if ml:
        _para(bt, False, ml, 23, MUTED, before_px=6, spc_em=0.12)


def _slide_contacts(slide, p, imgs):
    _bg(slide, DARK)
    # лого «гораздо больше» (правка 27.07)
    if not _logo(slide, LOGO_WHITE, 110, 52, 240):
        hdr = _txbox(slide, 110, 90, 1700, 44)
        _para(hdr, True, "STARGIFT", 26, "D0D0D0", spc_em=0.28)
    ht = _txbox(slide, 110, 324, 1400, 190)
    _para(ht, True, p.get("heading") or "Подарок для того, у кого есть всё",
          72, TXT_L, font=PRATA, line=1.08)
    it = _txbox(slide, 110, 512, 1350, 100)
    _para(it, True, p.get("intro") or ("Каждый экспонат сопровождается сертификатом "
          "подлинности. Оформление в раму на ваш вкус. Доставка по Москве "
          "и за её пределами."), 30, TXT_L2, line=1.5)
    if p.get("managerName"):
        _rect(slide, 110, 636, 700, 1, fill=LINE_D)
        mb = _txbox(slide, 110, 660, 1300, 160)
        _para(mb, True, "ВАШ ПЕРСОНАЛЬНЫЙ МЕНЕДЖЕР", 22, GOLD, spc_em=0.24)
        _para(mb, False, p["managerName"], 54, TXT_L, font=PRATA, before_px=12, line=1.05)
        if p.get("managerContact"):
            _para(mb, False, p["managerContact"], 30, "E6E6E6", before_px=10)
    # нижний ряд: только галереи, каждый адрес на своей строке (правка 27.07)
    _rect(slide, 110, 842, 1700, 1, fill=LINE_D)
    gb = _txbox(slide, 110, 864, 1700, 170)
    _para(gb, True, "ГАЛЕРЕИ", 22, MUTED, spc_em=0.22)
    _para(gb, False, "Галереи «Времена года» — Кутузовский проспект, 48", 28, TXT_L, before_px=12)
    _para(gb, False, "ТЦ «Гименей» — Большая Якиманка, 22", 28, TXT_L, before_px=8)
    _para(gb, False, "ТЦ «Дрим Хаус» — Барвиха, Жуковка, 85/1", 28, TXT_L, before_px=8)



def _scrim(slide, x, y, w, h, ang_deg=90, a0=0, a1=88, stops=None):
    """Прямоугольник с градиентом прозрачности чёрного (для читаемости подписей
    на полноэкранном кадре). По умолчанию 2 стопа (a0→a1); stops=[(permille,alpha%)]
    для многоточечного (редакционный разворот)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    sp.shadow.inherit = False
    sp.line.fill.background()
    spPr = sp.fill._xPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gs_lst = spPr.makeelement(qn("a:gsLst"), {})
    seq = stops if stops else [(0, a0), (100000, a1)]
    for pos, alpha in seq:
        gs = spPr.makeelement(qn("a:gs"), {"pos": str(int(pos))})
        clr = spPr.makeelement(qn("a:srgbClr"), {"val": "000000"})
        al = spPr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
        clr.append(al); gs.append(clr); gs_lst.append(gs)
    grad.append(gs_lst)
    lin = spPr.makeelement(qn("a:lin"), {"ang": str(int(ang_deg * 60000)), "scaled": "1"})
    grad.append(lin)
    ln = spPr.find(qn("a:ln"))
    spPr.insert(list(spPr).index(ln) if ln is not None else len(spPr), grad)
    return sp


def _slide_editorial(slide, imgs, kicker, title, paras):
    """Редакционный разворот: полноэкранный кадр + горизонтальный тёмный скрим
    слева + текстовая панель (кикер·золото, Prata-заголовок, абзацы) + врезки."""
    _bg(slide, DARK)
    _img_cover(slide, imgs.get("img0"), 0, 0, 1920, 1080)
    _scrim(slide, 0, 0, 1920, 1080, ang_deg=0,
           stops=[(0, 94), (30000, 86), (55000, 45), (75000, 10), (100000, 0)])
    tf = _txbox(slide, 100, 250, 1000, 620, "t")
    _para(tf, True, kicker, 26, GOLD, spc_em=0.26)
    for i, line in enumerate(title.split("\n")):
        _para(tf, False, line, 76, TXT_L, font=PRATA, line=1.04, before_px=(30 if i == 0 else 4))
    for j, t in enumerate(paras):
        _para(tf, False, t, 29, "E6E6E6", line=1.5, before_px=(36 if j == 0 else 20))
    # врезки-примеры снизу справа
    xs = 1420
    for src_key in ("img1", "img2"):
        if imgs.get(src_key):
            _img_cover(slide, imgs.get(src_key), xs, 700, 250, 320)
            _rect(slide, xs, 700, 250, 320, line="F5F5F5", line_w=1.5)
            xs += 272


def _full_caption(slide, lot, x, w, kicker, name_size, price_size, pad):
    """Подпись-оверлей на полноэкранном кадре: имя (Prata) + тип слева внизу,
    цена (золото) справа внизу; кикер — слева вверху."""
    if kicker:
        kt = _txbox(slide, x + pad, pad, w - 2 * pad, 40)
        _para(kt, True, kicker, 24, "F5F5F5", spc_em=0.26)
    person, subtitle, _ = _caption_parts(lot)
    top_txt, second_txt = (person, subtitle) if person else (subtitle, "")
    price_w = 190 if name_size < 48 else 280
    sub_size = 20 if name_size < 48 else 24
    nb = _txbox(slide, x + pad, 1080 - pad - 200, max(120, w - 2 * pad - price_w), 200, "b")
    first = True
    if top_txt:
        _para(nb, True, top_txt, name_size, "F5F5F5", font=PRATA, line=1.06); first = False
    if second_txt:
        _para(nb, first, second_txt, sub_size, "D9D9D9", before_px=(0 if first else 8))
    pr = _txbox(slide, x + w - pad - price_w, 1080 - pad - 60, price_w, 60, "b")
    _para(pr, True, _fmt_rub(lot["price"]), price_size, GOLD_HL, font=PRATA, align=PP_ALIGN.RIGHT)


def _slide_full_solo(slide, p, imgs):
    """n11 · Экспонат во весь экран."""
    _bg(slide, DARK)
    lots = p.get("lots") or []
    lot = lots[0] if lots else None
    if lot:
        _img_cover(slide, imgs.get(id(lot)), 0, 0, 1920, 1080)
    _scrim(slide, 0, 520, 1920, 560, a0=0, a1=90)              # нижний градиент
    _scrim(slide, 0, 0, 1920, 240, ang_deg=270, a0=0, a1=72)   # верхний градиент (для кикера)
    if lot:
        _full_caption(slide, lot, 0, 1920, p.get("kicker", ""), 78, 42, 90)


def _slide_full_strip(slide, p, imgs, n):
    """n13/n14 · Диптих/триптих во весь экран (кадры встык)."""
    _bg(slide, DARK)
    lots = (p.get("lots") or [])[:n]
    m = len(lots) or 1
    gap = 6
    pw = (1920 - gap * (m - 1)) / m
    for i, lot in enumerate(lots):
        x = i * (pw + gap)
        _img_cover(slide, imgs.get(id(lot)), x, 0, pw, 1080)
        _scrim(slide, x, 560, pw, 520, a0=0, a1=90)
        if i == 0:
            _scrim(slide, x, 0, pw, 200, ang_deg=270, a0=0, a1=72)
        _full_caption(slide, lot, x, pw, p.get("kicker", "") if i == 0 else "",
                      38 if n >= 3 else 52, 26 if n >= 3 else 34, 46)


def _mosaic_caption(slide, lot, x, w, y, name_size, price_size, price_bottom=None):
    """Подпись под фото в мозаике: имя+тип слева; цена (золото) справа —
    на НИЖНЕЙ линии ячейки (унификация 27.07), price_bottom = низ ячейки."""
    person, subtitle, _ = _caption_parts(lot)
    top_txt, second_txt = (person, subtitle) if person else (subtitle, "")
    price_w = 180
    nb = _txbox(slide, x, y, max(100, w - price_w), 80, "t")
    first = True
    if top_txt:
        _para(nb, True, top_txt, name_size, TXT_L, bold=True, line=1.18); first = False
    if second_txt:
        _para(nb, first, second_txt, name_size - 5, DESC_D, before_px=(0 if first else 4), line=1.2)
    py = (price_bottom - 44) if price_bottom else y
    pr = _txbox(slide, x + w - price_w, py, price_w, 44, "b" if price_bottom else "t")
    _para(pr, True, _fmt_rub(lot["price"]), price_size, GOLD_HL, font=PRATA, align=PP_ALIGN.RIGHT)


def _slide_mosaic(slide, p, imgs):
    """n15 · Мозаика 1+2: слева крупный кадр, справа два поменьше (тёмный фон)."""
    _bg(slide, DARK)
    lots = (p.get("lots") or [])[:3]
    kt = _txbox(slide, 90, 90, 1740, 40)
    _para(kt, True, p.get("kicker", ""), 26, GOLD, spc_em=0.26)
    ht = _txbox(slide, 90, 134, 1740, 90)
    _para(ht, True, (p.get("heading") or p.get("title") or "Экспонаты раздела"), 66, TXT_L, font=PRATA)
    top, bottom = 250, 1000
    lw, gap = 1044, 44
    rw = 1740 - lw - gap
    if lots:
        big_ph = (bottom - top) - 96
        d = _img_fill_h(slide, imgs.get(id(lots[0])), 90, lw, top, big_ph)
        cy = (d[1] + d[3]) if d else (top + big_ph)
        _mosaic_caption(slide, lots[0], 90, lw, cy + 16, 26, 28, price_bottom=1000)
    rx = 90 + lw + gap
    ch = (bottom - top - 40) / 2.0
    for i, lot in enumerate(lots[1:3]):
        cy0 = top + i * (ch + 40)
        small_ph = ch - 80
        d = _img_fill_h(slide, imgs.get(id(lot)), rx, rw, cy0, small_ph)
        cy = (d[1] + d[3]) if d else (cy0 + small_ph)
        _mosaic_caption(slide, lot, rx, rw, cy + 12, 22, 23, price_bottom=cy0 + ch)


def _slide_clean_divider(slide, p, imgs):
    """q12 · Чистый разделитель темы: типографский шмуцтитул без фото (тёмный)."""
    _bg(slide, DARK)
    if p.get("roman"):
        rt = _txbox(slide, 120, 300, 1400, 60)
        _para(rt, True, p["roman"], 40, MUTED, font=PRATA)
    ht = _txbox(slide, 120, 370, 1680, 260)
    _para(ht, True, p.get("title", ""), 104, TXT_L, font=PRATA, line=1.02)
    y = 630
    if p.get("subtitle"):
        it = _txbox(slide, 120, y, 1150, 160)
        _para(it, True, p["subtitle"], 32, TXT_L2, line=1.5)
        y += 170
    cnt = int(p.get("count", 0) or 0)
    if cnt:
        tail = "%d %s" % (cnt, _plural(cnt, "ЭКСПОНАТ", "ЭКСПОНАТА", "ЭКСПОНАТОВ"))
        if p.get("min"):
            tail += ", ОТ " + _fmt_rub(p["min"])
        ct = _txbox(slide, 120, y, 1400, 44)
        _para(ct, True, tail, 26, GOLD, spc_em=0.22)


_FRAMING_TXT = ("Оформление на Ваш вкус, подарочная упаковка, доставка по Москве "
                "и помощь в организации доставки по миру входят в стоимость. "
                "Для журналов, сценариев и книг — рамы с механизмом-шкатулкой: "
                "экспонат можно извлечь, полистать и вернуть обратно.")


def _slide_framing(slide, p, imgs):
    """Подарочное оформление. light: чистый светлый слайд — текст + ряд из 3 фото;
    иначе тёмный редакционный разворот (каталожный режим)."""
    txt = p.get("text") or _FRAMING_TXT
    if p.get("light"):
        _bg(slide, PAPER)
        kt = _txbox(slide, 110, 84, 1700, 40)
        _para(kt, True, "СЕРВИС STARGIFT", 26, MUTED, spc_em=0.26)
        ht = _txbox(slide, 110, 134, 1700, 100)
        _para(ht, True, "Подарочное оформление", 66, TXT_D, font=PRATA)
        tt = _txbox(slide, 110, 244, 1560, 140)
        _para(tt, True, txt, 29, "3A3A3A", line=1.5)
        # ряд из трёх фото ровной высоты
        keys = [k for k in ("img0", "img1", "img2") if imgs.get(k)]
        n = len(keys) or 1
        gap = 56
        cw = (1700 - gap * (n - 1)) / n
        for i, k in enumerate(keys):
            _img_contain(slide, imgs.get(k), 110 + i * (cw + gap), 430, cw, 560)
        return
    _slide_editorial(slide, imgs, "СЕРВИС STARGIFT", "Подарочное\nоформление", [txt])


def _slide_delivery(slide, p, imgs):
    """Доставка и подарочный сертификат. light: текст слева + крупный сертификат
    справа на белом; иначе тёмный редакционный разворот (каталожный режим)."""
    texts = p.get("texts") or [
        "Если нужного экспоната нет в коллекции, а подарок нужен оперативно — мы найдём "
        "его для вас по запросу. Доставка из США и Европы в Москву занимает около 2 месяцев.",
        "В день торжества одариваемому вручается подарочный сертификат. По прибытии "
        "в Москву мы оформим экспонат в раму по вашим предпочтениям за 1–2 дня."]
    if p.get("light"):
        _bg(slide, PAPER)
        tf = _txbox(slide, 110, 200, 900, 700, "m")
        _para(tf, True, "ДОСТАВКА И ПОДАРОЧНЫЙ СЕРТИФИКАТ", 26, MUTED, spc_em=0.26)
        _para(tf, False, "Подарок готов", 66, TXT_D, font=PRATA, before_px=30, line=1.08)
        _para(tf, False, "в день торжества", 66, TXT_D, font=PRATA, before_px=4, line=1.08)
        for t in texts:
            _para(tf, False, t, 29, "3A3A3A", before_px=28, line=1.55)
        _img_contain(slide, imgs.get("img0"), 1080, 120, 730, 840)
        return
    _slide_editorial(slide, imgs, "ДОСТАВКА И ПОДАРОЧНЫЙ СЕРТИФИКАТ",
                     "Подарок готов\nв день торжества", texts)


# ──────────────────────────────────────────────────────────────────────────
#  НОВЫЕ ФОРМАТЫ (зеркало новых TPL из constructor.html):
#  macro-sign · detail-shot · mood-strip · manifest · mood-lots ·
#  custom-1 «Цена как арт-объект» · custom-2 «Две эпохи»
# ──────────────────────────────────────────────────────────────────────────
def _img_crop_zoom(slide, path, x, y, w, h, fx=0.5, fy=0.5, zoom=1.0):
    """Кроп-зум фото ровно в область w×h — аналог веб-связки
    object-fit:cover + object-position:fx fy + transform:scale(zoom).
    В PPTX увеличения нет, поэтому видимая часть источника задаётся
    crop-атрибутами картинки: чем больше zoom, тем меньше окно кропа.
    fx/fy — точка фокуса в долях (0.5/0.62 = центр по X, ниже центра по Y)."""
    if not path:
        return None
    try:
        pic = slide.shapes.add_picture(str(path), E(x), E(y))
        a = pic.width / float(pic.height)          # природный аспект источника
        t = float(w) / float(h)                    # аспект области
        z = max(1.0, float(zoom))
        if a > t:                                  # источник шире — по высоте влезает целиком
            fh, fw = 1.0 / z, (t / a) / z
        else:                                      # источник выше области
            fw, fh = 1.0 / z, (a / t) / z
        fw, fh = min(1.0, max(0.02, fw)), min(1.0, max(0.02, fh))
        fx = min(1.0, max(0.0, fx)); fy = min(1.0, max(0.0, fy))
        pic.crop_left = fx * (1 - fw)
        pic.crop_right = (1 - fx) * (1 - fw)
        pic.crop_top = fy * (1 - fh)
        pic.crop_bottom = (1 - fy) * (1 - fh)
        pic.left, pic.top = E(x), E(y)
        pic.width, pic.height = E(w), E(h)
        return pic
    except Exception:
        return None


def _run(p, text, size_px, color, font=GOLOS, bold=False, spc_em=None, upper=False):
    """Дописать ещё один run в УЖЕ созданный абзац — для строк, собранных из
    разноцветных фрагментов («Имя · тип» одной строкой)."""
    r = p.add_run()
    r.text = (text or "").upper() if upper else (text or "")
    r.font.name = font
    r.font.size = PT(size_px)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    if spc_em:
        _set_spc(r, spc_em, size_px)
    return r


def _cpl58(w, size_px):
    """Консервативная оценка «символов в строке» (шире, чем _cap_lines с 0.52).
    Нужна там, где подпись стоит вплотную к фиксированной строке цены: лучше
    переоценить длину строки и подрезать раньше, чем получить наезд."""
    return max(1, int(w / (size_px * 0.58)))


def _nlines58(txt, size_px, w, maxl):
    import math
    if not txt:
        return 0
    return min(maxl, max(1, int(math.ceil(len(txt) / float(_cpl58(w, size_px))))))


_LSP = 1.22   # PowerPoint считает line_spacing от «одинарного» интервала шрифта,
              # а не от кегля: реальная высота строки ≈ кегль × множитель × _LSP


def _lh(size_px, mult=1.0):
    """Реальная высота строки в px для расчёта зон (см. _LSP)."""
    return size_px * mult * _LSP


def _clamp_lines(txt, size_px, w, nlines):
    """Обрезать текст многоточием под nlines строк (аналог -webkit-line-clamp)."""
    if not txt or nlines <= 0:
        return ""
    budget = nlines * _cpl58(w, size_px)
    if len(txt) <= budget:
        return txt
    return txt[:budget - 1].rstrip() + "…"


def _first_lot(p):
    lots = p.get("lots") or ([p["lot"]] if p.get("lot") else [])
    return lots[0] if lots else {}


def _two_levels(lot):
    """(крупный уровень, второй уровень) с фолбэком на полное имя лота."""
    person, subtitle, full = _caption_parts(lot)
    if person:
        return person, subtitle
    return (subtitle or full), ""


def _slide_macro_sign(slide, p, imgs):
    """n07 · Макро: подпись во весь экран. Кадр экспоната сильно увеличен
    (зум 2.0, фокус 50%/62% — зона автографа), снизу градиентный скрим и одна
    тонкая строка: имя · тип слева, цена золотом справа. Больше ничего."""
    _bg(slide, DARK)
    lot = _first_lot(p)
    _img_crop_zoom(slide, imgs.get(id(lot)), 0, 0, 1920, 1080, 0.5, 0.62, 1.45)
    _scrim(slide, 0, 0, 1920, 1080, ang_deg=90,
           stops=[(0, 25), (22000, 0), (46000, 0), (100000, 88)])
    top_txt, second_txt = _two_levels(lot)
    # блок подписи прижат к низу (как bottom:74 в вебе): длинная строка переносится
    # на вторую строку и уезжает ВВЕРХ вместе с линией, а не за нижний край слайда
    joined = top_txt + (" · " if top_txt and second_txt else "") + second_txt
    nlines = _nlines58(joined, 30, 1300, 2)
    txt_top = 1006 - max(1, nlines) * _lh(30, 1.32)
    _rect(slide, 90, txt_top - 26, 1740, 1, fill=TXT_L, alpha=28)
    tf = _txbox(slide, 90, txt_top, 1300, 90, "t")
    pp, _r = _para(tf, True, top_txt, 30, TXT_L, bold=True, line=1.32)
    if top_txt and second_txt:
        _run(pp, " · ", 30, MUTED)
    if second_txt:
        _run(pp, second_txt, 30, TXT_L2)
    pr = _txbox(slide, 1430, txt_top - 6, 400, 70, "t")
    _para(pr, True, _fmt_rub(lot.get("price", 0)), 36, GOLD_HL, font=PRATA,
          align=PP_ALIGN.RIGHT)


def _slide_detail_shot(slide, p, imgs):
    """n12 · Общий план + детали. Слева крупный кадр экспоната целиком и подпись
    с ценой под ним, справа два квадрата 340×340 — кроп-зумы ТОГО ЖЕ файла
    («АВТОГРАФ» 50%/62% зум 2.2 и «ДЕТАЛЬ» 50%/26% зум 1.7). Тёмный фон."""
    _bg(slide, DARK)
    lot = _first_lot(p)
    src = imgs.get(id(lot))
    kt = _txbox(slide, 90, 88, 1740, 80)
    _para(kt, True, p.get("kicker", ""), 26, GOLD, spc_em=0.26)
    dy = 34 if len(p.get("kicker", "")) * 26 * 0.84 > 1740 else 0
    ht = _txbox(slide, 90, 132 + dy, 1740, 96)
    _para(ht, True, (p.get("heading") or p.get("title") or ""), 64, TXT_L, font=PRATA)
    top, bottom = 252 + dy, 1002
    lw, gap = 1348, 52
    rx, rw = 90 + lw + gap, 340
    top_txt, second_txt = _two_levels(lot)
    cert = (lot.get("cert") or "").strip()
    cert = cert if (cert and cert != "Provenance") else ""
    # подпись прижата к низу (как в вебе) и растёт ВВЕРХ, отбирая высоту у фото
    cap_w = lw - 400
    cap_h = _nlines58(top_txt, 32, cap_w, 2) * _lh(32, 1.22)
    if second_txt:
        cap_h += 7 + _nlines58(second_txt, 24, cap_w, 2) * _lh(24, 1.3)
    if cert:
        cap_h += 7 + _lh(21, 1.4)
    rule_y = int(bottom - cap_h - 22)
    _img_contain(slide, src, 90, top, lw, max(80, (rule_y - 26) - top))
    _rect(slide, 90, rule_y, lw, 1, fill=LINE_D)
    nb = _txbox(slide, 90, rule_y + 22, cap_w, 160, "t")
    first = True
    if top_txt:
        _para(nb, True, top_txt, 32, TXT_L, bold=True, line=1.22); first = False
    if second_txt:
        _para(nb, first, second_txt, 24, DESC_D,
              before_px=(0 if first else 7), line=1.3); first = False
    if cert:
        _para(nb, first, cert, 21, MUTED, before_px=(0 if first else 7),
              spc_em=0.14, upper=True)
    pb = _txbox(slide, 90 + lw - 400, rule_y + 22, 400, 60, "t")
    _para(pb, True, _fmt_rub(lot.get("price", 0)), 34, GOLD_HL, font=PRATA,
          align=PP_ALIGN.RIGHT)
    # два квадратных кроп-зума того же кадра, по центру правой колонки
    crops = [("АВТОГРАФ", 0.5, 0.62, 2.2), ("ДЕТАЛЬ", 0.5, 0.26, 1.7)]
    y0 = top + ((bottom - top) - (2 * rw + 40)) / 2.0
    for i, (label, fx, fy, z) in enumerate(crops):
        cy = y0 + i * (rw + 40)
        _rect(slide, rx, cy, rw, rw, fill=CARD_D)
        _img_crop_zoom(slide, src, rx, cy, rw, rw, fx, fy, z)
        _scrim(slide, rx, cy + rw - 72, rw, 72, ang_deg=90, a0=0, a1=82)
        _rect(slide, rx, cy, rw, rw, fill=None, line=LINE_D)
        lt = _txbox(slide, rx + 18, cy + rw - 44, rw - 36, 34, "t")
        _para(lt, True, label, 19, GOLD, spc_em=0.2)


def _slide_mood_strip(slide, p, imgs):
    """k02 · Три кадра настроения: архивные фото в тонких тёмных рамах, кикер и
    заголовок сверху, слева внизу пояснение, справа внизу источник. Без лотов."""
    _bg(slide, DARK)
    kt = _txbox(slide, 90, 88, 1740, 80)
    _para(kt, True, p.get("kicker", ""), 26, GOLD, spc_em=0.26)
    dy = 34 if len(p.get("kicker", "")) * 26 * 0.84 > 1740 else 0
    ht = _txbox(slide, 90, 132 + dy, 1740, 100)
    _para(ht, True, (p.get("heading") or p.get("title") or ""), 68, TXT_L, font=PRATA)
    keys = [i for i, s in enumerate(p.get("imgs") or []) if s][:3]
    top, grid_b = 256 + dy, 924
    n = len(keys) or 1
    gap, pad = 36, 14
    cw = (1740 - gap * (n - 1)) / float(n)
    for j, i in enumerate(keys):
        cx = 90 + j * (cw + gap)
        _rect(slide, cx, top, cw, grid_b - top, fill=CARD_D, line=LINE_D)
        _img_cover(slide, imgs.get("img%d" % i), cx + pad, top + pad,
                   cw - 2 * pad, (grid_b - top) - 2 * pad)
    if p.get("note"):
        nt = _txbox(slide, 90, 954, 1180, 70, "t")
        _para(nt, True, p["note"], 23, MUTED, line=1.4)
    if p.get("credit"):
        ct = _txbox(slide, 1290, 956, 540, 60, "t")
        _para(ct, True, p["credit"], 19, MUTED, spc_em=0.14,
              align=PP_ALIGN.RIGHT, upper=True)


def _slide_manifest(slide, p, imgs):
    """n01 · Манифест с цитатой: чёрный слайд, кикер золотом, крупная цитата Prata
    в кавычках-ёлочках, атрибуция «— автор», внизу над линией строка фактов капсом."""
    _bg(slide, DARK)
    foot = p.get("footnote")
    if isinstance(foot, (list, tuple)):
        foot = " · ".join([x for x in foot if x])
    foot = (foot or "").strip()
    quote = (p.get("quote") or "").strip()
    box_h = 780 if foot else 880
    tf = _txbox(slide, 120, 100, 1560, box_h, "m")
    first = True
    if p.get("kicker"):
        _para(tf, True, p["kicker"], 26, GOLD, spc_em=0.26); first = False
    # кегль цитаты снижаем на длинных текстах, чтобы влезала без обрезки
    qs = 68
    if len(quote) > 210:
        qs = 52
    elif len(quote) > 150:
        qs = 58
    _para(tf, first, "«%s»" % quote, qs, TXT_L, font=PRATA,
          before_px=(0 if first else 52), line=1.2); first = False
    if p.get("author"):
        _para(tf, first, "— " + p["author"], 29, DESC_D,
              before_px=(0 if first else 48), spc_em=0.1)
    if foot:
        _rect(slide, 120, 900, 1560, 1, fill=LINE_D)
        ft = _txbox(slide, 120, 934, 1560, 120, "t")
        _para(ft, True, foot, 24, MUTED, spc_em=0.22, line=1.6, upper=True)


def _slide_mood_lots(slide, p, imgs):
    """q10 · Кадр-метафора + строка экспонатов: полноэкранный тематический кадр
    под градиентом, сверху слева кикер и заголовок, внизу три мини-блока
    (фото 120×120 + имя/тип + цена золотом), справа внизу источник кадра."""
    _bg(slide, DARK)
    _img_cover(slide, imgs.get("bg"), 0, 0, 1920, 1080)
    _scrim(slide, 0, 0, 1920, 1080, ang_deg=90,
           stops=[(0, 82), (32000, 28), (62000, 62), (100000, 95)])
    kt = _txbox(slide, 90, 88, 1250, 80)
    _para(kt, True, p.get("kicker", ""), 26, GOLD, spc_em=0.26)
    ht = _txbox(slide, 90, 146 + (34 if len(p.get("kicker", "")) * 26 * 0.84 > 1250 else 0),
                1250, 190)
    _para(ht, True, (p.get("heading") or p.get("title") or ""), 76, TXT_L,
          font=PRATA, line=1.04)
    lots = (p.get("lots") or [])[:3]
    n = len(lots) or 1
    gap, ph = 44, 120
    cw = (1740 - gap * (n - 1)) / float(n)
    _rect(slide, 90, 812, 1740, 1, fill=TXT_L, alpha=22)
    row_y, price_y, zone_top = 842, 962, 830
    for i, lot in enumerate(lots):
        cx = 90 + i * (cw + gap)
        _rect(slide, cx, row_y, ph, ph, fill=CARD_D)
        _img_cover(slide, imgs.get(id(lot)), cx, row_y, ph, ph)
        _rect(slide, cx, row_y, ph, ph, fill=None, line="3F3F3F")
        tx, tw = cx + ph + 22, cw - ph - 22
        top_txt, second_txt = _two_levels(lot)
        # имя не режем (до 2 строк), подзаголовок подрезаем под остаток —
        # аналог -webkit-line-clamp:2 в вебе; цена у всех блоков на общей Y
        nl = _nlines58(top_txt, 25, tw, 2)
        avail = (price_y - 8 - zone_top) - nl * _lh(25, 1.24) - (8 if second_txt else 0)
        sl = max(0, min(2, int(avail / _lh(20, 1.28))))
        second_txt = _clamp_lines(second_txt, 20, tw, sl)
        tf = _txbox(slide, tx, zone_top, tw, price_y - 8 - zone_top, "t")
        first = True
        if top_txt:
            _para(tf, True, top_txt, 25, TXT_L, bold=True, line=1.24); first = False
        if second_txt:
            _para(tf, first, second_txt, 20, TXT_L2,
                  before_px=(0 if first else 8), line=1.28)
        pb = _txbox(slide, tx, price_y, tw, 50, "t")
        _para(pb, True, _fmt_rub(lot.get("price", 0)), 25, GOLD_HL, font=PRATA)
    if p.get("credit"):
        cr = _txbox(slide, 990, 1024, 840, 40, "t")
        _para(cr, True, p["credit"], 18, MUTED, spc_em=0.14,
              align=PP_ALIGN.RIGHT, upper=True)


def _slide_price_art(slide, p, imgs):
    """АВТОРСКИЙ 1 · «Цена как арт-объект». Гигантская цена Prata по центру-верху
    (кегль от длины строки, чтобы влезала в 1740px), под ней фото экспоната,
    внизу имя+тип слева и мелкая спецификация справа. В PPTX контурного текста
    нет — вместо обводки заливка GOLD_HL при крупном кегле."""
    _bg(slide, DARK)
    lot = _first_lot(p)
    kt = _txbox(slide, 90, 80, 1740, 40)
    _para(kt, True, p.get("kicker") or "ЭКСПОНАТ", 26, GOLD, spc_em=0.26)
    txt = _fmt_rub(lot.get("price", 0))
    fs = int(min(196, 1740 / (max(1, len(txt)) * 0.62)))
    pt = _txbox(slide, 90, 130, 1740, fs * 1.2, "t")
    pt.word_wrap = False
    _para(pt, True, txt, fs, GOLD_HL, font=PRATA, line=1.06, align=PP_ALIGN.CENTER)
    rule_y = 895
    _rect(slide, 90, rule_y, 1740, 1, fill=LINE_D)
    ph_top = 130 + int(fs * 1.16)
    _img_contain(slide, imgs.get(id(lot)), 403, ph_top, 1114,
                 max(80, (rule_y - 26) - ph_top))
    top_txt, second_txt = _two_levels(lot)
    nb = _txbox(slide, 90, rule_y + 26, 1140, 120, "t")
    first = True
    if top_txt:
        _para(nb, True, top_txt, 42, TXT_L, font=PRATA, line=1.1); first = False
    if second_txt:
        _para(nb, first, second_txt, 24, TXT_L2,
              before_px=(0 if first else 8), line=1.3)
    note = p.get("note") or ("в наличии в Москве" if lot.get("instock") else "под заказ")
    nt = _txbox(slide, 1280, rule_y + 32, 550, 60, "t")
    _para(nt, True, note, 21, MUTED, spc_em=0.2, align=PP_ALIGN.RIGHT, upper=True)


def _slide_two_eras(slide, p, imgs):
    """АВТОРСКИЙ 2 · «Две эпохи»: разворот-сравнение. Слева светлая карточка
    (эпоха-первоисточник), справа тёмная (сегодня), между ними тонкая линия
    с золотой точкой. Сверху у каждой — капс-лейбл эпохи, внутри фото и подпись
    с ценой. Справа сверху — охват периода (span)."""
    _bg(slide, DARK)
    span = (p.get("span") or "").strip()
    kicker = p.get("kicker", "")
    kw = 1000 if span else 1520
    kt = _txbox(slide, 90, 84, kw, 80)
    _para(kt, True, kicker, 26, GOLD, spc_em=0.26)
    if span:
        st = _txbox(slide, 1130, 86, 700, 40, "t")
        _para(st, True, span, 22, MUTED, spc_em=0.2,
              align=PP_ALIGN.RIGHT, upper=True)
    # длинный кикер переносится на вторую строку — заголовок уезжает ниже
    kl = 1 if len(kicker) * 26 * 0.84 <= kw else 2
    ht = _txbox(slide, 90, 130 + (34 if kl > 1 else 0), 1740, 96)
    _para(ht, True, (p.get("heading") or p.get("title") or ""), 64, TXT_L, font=PRATA)
    lots = (p.get("lots") or [])[:2]
    # подписи-панелей ставим ТОЛЬКО если переданы реальные эпохи (как в вебе):
    # годов в каталоге нет, «ТОГДА/СЕЙЧАС» на двух современниках — выдумка
    eras = [p.get("eraLeft") or "", p.get("eraRight") or ""]
    top, bottom = 250, 1006
    pw = (1740 - 64) / 2.0
    # разделитель: тонкая линия по центру желоба + золотая точка
    _rect(slide, 960, top, 1, bottom - top, fill=LINE_D)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(953), E((top + bottom) / 2.0 - 7),
                                 E(14), E(14))
    dot.shadow.inherit = False
    dot.fill.solid(); dot.fill.fore_color.rgb = C(GOLD)
    dot.line.fill.background()
    for i, lot in enumerate(lots):
        dark = (i == 1)
        px = 90 if i == 0 else (90 + pw + 64)
        _rect(slide, px, top, pw, bottom - top,
              fill=(CARD_D if dark else CARD_L), line=(LINE_D if dark else LINE_L))
        name_c = TXT_L if dark else TXT_D
        sub_c = DESC_D if dark else DESC_L
        era_c = GOLD if dark else GOLD_L
        price_c = GOLD_HL if dark else GOLD_L
        et = _txbox(slide, px + 44, top + 46, pw - 88, 40, "t")
        _para(et, True, eras[i], 24, era_c, spc_em=0.24, upper=True)
        _img_contain(slide, imgs.get(id(lot)), px + 44, top + 106, pw - 88, 396)
        top_txt, second_txt = _two_levels(lot)
        cw2 = pw - 88
        cap_top, price_y = 768, bottom - 44 - 45
        # имя целиком (до 2 строк), подзаголовок подрезаем под остаток зоны —
        # цена у обеих карточек стоит на одной фиксированной Y
        nl = _nlines58(top_txt, 30, cw2, 2)
        avail = (price_y - 8 - cap_top) - nl * _lh(30, 1.22) - (8 if second_txt else 0)
        sl = max(0, min(2, int(avail / _lh(23, 1.3))))
        second_txt = _clamp_lines(second_txt, 23, cw2, sl)
        cb = _txbox(slide, px + 44, cap_top, cw2, price_y - 8 - cap_top, "t")
        first = True
        if top_txt:
            _para(cb, True, top_txt, 30, name_c, bold=True, line=1.22); first = False
        if second_txt:
            _para(cb, first, second_txt, 23, sub_c,
                  before_px=(0 if first else 8), line=1.3)
        prt = _txbox(slide, px + 44, price_y, cw2, 60, "t")
        _para(prt, True, _fmt_rub(lot.get("price", 0)), 30, price_c, font=PRATA)


# тип слайда → (рендер, ожидает ли лоты)
_RENDER = {
    "cover": _slide_cover,
    "divider": _slide_divider,
    "contents": _slide_contents,
    "premium-solo": _slide_premium_solo,
    "grid4": _slide_grid4,
    "pair-dark": _slide_pair_dark,
    "pair-light": _slide_pair_light,
    "cards-side": _slide_cards_side,
    "trio": lambda s, p, im: _slide_trio(s, p, im, False),
    "trio-dark": lambda s, p, im: _slide_trio(s, p, im, True),
    "duo-clean": lambda s, p, im: _slide_duo_clean(s, p, im, False),
    "duo-clean-dark": lambda s, p, im: _slide_duo_clean(s, p, im, True),
    "full-solo": _slide_full_solo,
    "full-duo": lambda s, p, im: _slide_full_strip(s, p, im, 2),
    "full-trio": lambda s, p, im: _slide_full_strip(s, p, im, 3),
    "mosaic": _slide_mosaic,
    "clean-divider": _slide_clean_divider,
    "instock": _slide_instock,
    "index": _slide_index,
    "quote": _slide_quote,
    "framing": _slide_framing,
    "delivery": _slide_delivery,
    "contacts": _slide_contacts,
    # новые форматы (портал шаблонов + авторские)
    "macro-sign": _slide_macro_sign,
    "detail-shot": _slide_detail_shot,
    "mood-strip": _slide_mood_strip,
    "manifest": _slide_manifest,
    "mood-lots": _slide_mood_lots,
    "custom-1": _slide_price_art,
    "custom-2": _slide_two_eras,
}


def _collect_photos(slide_data, workdir, counter):
    """Скачать все фото слайда, вернуть карту {ключ → путь}. Ключи:
    'cover'/'bg' для фонов, id(lot) для каждого лота."""
    imgs = {}
    params = slide_data.get("params", {}) or {}
    typ = slide_data.get("type")
    if typ == "cover" and params.get("cover"):
        imgs["cover"] = _fetch(params["cover"], workdir, next(counter))
    if typ == "divider" and params.get("bg"):
        imgs["bg"] = _fetch(params["bg"], workdir, next(counter))
    # mood-lots: полноэкранный кадр-метафора лежит в params.bg (как у divider)
    if typ == "mood-lots" and params.get("bg"):
        imgs["bg"] = _fetch(params["bg"], workdir, next(counter))
    lots = list(slide_data.get("lots", []) or [])
    if params.get("lot"):
        lots.append(params["lot"])
    for lot in lots:
        if lot.get("img"):
            imgs[id(lot)] = _fetch(lot["img"], workdir, next(counter))
    # framing/delivery несут список фото в params.imgs (не лоты) → ключи img0/img1/…
    for i, src in enumerate(params.get("imgs") or []):
        if src:
            imgs["img%d" % i] = _fetch(src, workdir, next(counter))
    return imgs


def _merge_params(slide_data):
    """Свести params + lots/lot в единый словарь p, ожидаемый рендерами."""
    p = dict(slide_data.get("params", {}) or {})
    if "lots" in slide_data:
        p["lots"] = slide_data["lots"]
    return p


def build_pptx(data, out_path):
    """Собрать .pptx из data. Возвращает Path готового файла."""
    out_path = Path(out_path)
    _FETCH_CACHE.clear()          # кэш живёт в пределах одной сборки
    prs = Presentation()
    prs.slide_width = E(1920)
    prs.slide_height = E(1080)
    blank = prs.slide_layouts[6]

    def _counter():
        i = 0
        while True:
            yield i
            i += 1

    with tempfile.TemporaryDirectory() as td:
        workdir = Path(td)
        counter = _counter()
        for sd in data.get("slides", []):
            typ = sd.get("type")
            render = _RENDER.get(typ)
            if not render:
                continue
            slide = prs.slides.add_slide(blank)
            try:
                imgs = _collect_photos(sd, workdir, counter)
                render(slide, _merge_params(sd), imgs)
            except Exception:
                # один сбойный слайд не должен ронять всю сборку
                pass

        prs.save(str(out_path))
    return out_path
