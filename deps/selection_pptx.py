"""Редактируемая версия подборки — PowerPoint (.pptx), ЗЕРКАЛО selection_pdf.

Раскладка, кегли, веса шрифта, отступы и варианты фото-колонки повторяют
PDF-эталон 1:1 (сверка Вашика 28.07): 4:3 (1920×1440 px холста → слайд 10×7.5"),
фото-колонка слева (варианты A/B/C как в PDF), текст справа с PDF-отступами
(34/42/18/42/42/46 px), скидочные цены, размеры, сертификат, ссылка тёмным
в потоке, финальный слайд с вертикально отцентрованной колонкой.
Все тексты — отдельными текстбоксами (фидбек Насти 28.07), менеджер правит
в PowerPoint/Keynote.
"""
import tempfile
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GOLD = RGBColor(0x9A, 0x82, 0x53)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
BLACK = RGBColor(0, 0, 0)

# px(1920-холст) → дюймы слайда 10"
def _in(px: float) -> Inches:
    return Inches(px / 192.0)

ASSETS = Path(__file__).parent / "selection_pdf_assets"  # тот же каталог, что у PDF


def _fetch(url: str, workdir: Path, idx: int):
    """Скачать фото во временный файл; None при неудаче (слайд соберётся без него)."""
    if not url:
        return None
    try:
        ext = ".jpg"
        low = url.lower()
        for e in (".png", ".webp", ".jpeg"):
            if low.endswith(e):
                ext = e
        # Локальный путь (фото из чата) — читаем с диска, иначе качаем.
        if url.startswith("/") and Path(url).exists():
            src = Path(url)
        else:
            # У части карточек в каталоге ОТНОСИТЕЛЬНЫЕ пути (/media/…) —
            # без домена фото молча пропадало (инцидент 28.07, Мэджик Джонсон).
            if url.startswith("/"):
                url = "https://stargift.ru" + url
            from urllib.parse import urlsplit, urlunsplit, quote
            p = urlsplit(url)
            url = urlunsplit((p.scheme, p.netloc, quote(p.path, safe="/%"), p.query, ""))
            src = workdir / f"pp_raw_{idx}{ext}"
            req = urllib.request.Request(url, headers={"User-Agent": "stargift-doc-bot"})
            try:
                with urllib.request.urlopen(req, timeout=30) as r:
                    src.write_bytes(r.read())
            except Exception:
                with urllib.request.urlopen(req, timeout=60) as r:  # один ретрай
                    src.write_bytes(r.read())
        # Ужимаем фото — иначе Keynote/pptx раздувается до 20+ МБ (полное
        # разрешение) и отправка падает с 413, из-за чего Док дробил презентацию
        # (инцидент 24.07). 1600px/q88 — щадящий баланс после жалобы Насти 28.07
        # на пережатые фото (26 слайдов всё равно укладываются в лимит).
        from PIL import Image
        out = workdir / f"pp_{idx}.jpg"
        im = Image.open(src).convert("RGB")
        if max(im.size) > 1600:
            im.thumbnail((1600, 1600))
        im.save(out, "JPEG", quality=88)
        return out
    except Exception:
        return None


def _cover_crop(path, workdir: Path, tag: str, w_px: float, h_px: float):
    """Центральный кроп под пропорцию рамки — зеркало CSS object-fit: cover."""
    try:
        from PIL import Image
        im = Image.open(path)
        tw, th = im.size
        target = w_px / h_px
        cur = tw / th
        if cur > target:  # шире рамки — режем бока
            nw = int(th * target)
            box = ((tw - nw) // 2, 0, (tw - nw) // 2 + nw, th)
        else:             # выше рамки — режем верх/низ
            nh = int(tw / target)
            box = (0, (th - nh) // 2, tw, (th - nh) // 2 + nh)
        out = workdir / f"cov_{tag}.jpg"
        im.crop(box).convert("RGB").save(out, "JPEG", quality=88)
        return out
    except Exception:
        return path


def _gold_line(shape):
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)  # 3px холста 1920 ≈ 1.5pt


def _pic_cover(slide, path, left, top, width, height, workdir, tag):
    """Фото, заполняющее рамку с обрезкой (PDF .photo-bottom: object-fit cover)."""
    if not path:
        return
    try:
        cropped = _cover_crop(path, workdir, tag, width, height)
        pic = slide.shapes.add_picture(str(cropped), _in(left), _in(top),
                                       width=_in(width), height=_in(height))
        _gold_line(pic)
    except Exception:
        pass


def _pic_contain_framed(slide, path, left, top, width, height):
    """Рамка-контейнер с золотым бортом + фото целиком внутри (PDF .photo-top
    .frame: рамка на контейнере, белые поля внутри при несовпадении пропорций)."""
    if not path:
        return
    try:
        from pptx.enum.shapes import MSO_SHAPE
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _in(left), _in(top),
                                      _in(width), _in(height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _gold_line(rect)
        rect.shadow.inherit = False
        pic = slide.shapes.add_picture(str(path), _in(left), _in(top), height=_in(height - 8))
        max_w = _in(width - 8)
        if pic.width > max_w:
            ratio = max_w / pic.width
            pic.width = max_w
            pic.height = Emu(int(pic.height * ratio))
        pic.left = Emu(int(_in(left) + (_in(width) - pic.width) / 2))
        pic.top = Emu(int(_in(top) + (_in(height) - pic.height) / 2))
    except Exception:
        pass


def _pic_contain_selfborder(slide, path, left, top, width, height):
    """Фото целиком, рамка на самом фото (PDF .photo-full — вариант B)."""
    if not path:
        return
    try:
        pic = slide.shapes.add_picture(str(path), _in(left), _in(top), height=_in(height))
        if pic.width > _in(width):
            ratio = _in(width) / pic.width
            pic.width = _in(width)
            pic.height = Emu(int(pic.height * ratio))
        pic.left = Emu(int(_in(left) + (_in(width) - pic.width) / 2))
        pic.top = Emu(int(_in(top) + (_in(height) - pic.height) / 2))
        _gold_line(pic)
    except Exception:
        pass


# Семейства шрифта: в наборе Proxima Nova регуляр живёт ОТДЕЛЬНЫМ семейством
# «Proxima Nova Rg», а лёгкое начертание (вес PDF-текста) — «Proxima Nova Light»;
# в семействе «Proxima Nova» только Light/Bold/ExtraBold/Black, и запрос
# регуляра из него рендерится Bold'ом («всё жирное», инцидент 28.07).
FONT = "Proxima Nova"          # для жирных надписей (имя, скидочная цена)
FONT_RG = "Proxima Nova Rg"    # весь остальной текст (единственное семейство
                               # с регулярным весом; «Proxima Nova Light» как
                               # СЕМЕЙСТВО в наборе не существует — подменялся
                               # системным шрифтом, жалоба Вашика 28.07)


def _style(r, size, bold=False, color=DARK, underline=False):
    r.font.name = FONT if bold else FONT_RG
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.underline = underline
    r.font.color.rgb = color


def _txbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    # нулевые внутренние поля — иначе текст съезжает относительно PDF-сетки
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    try:
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
        tf.auto_size = MSO_AUTO_SIZE.NONE  # единый кегль считается заранее — без локальных ужатий
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    return tf


def _est_lines(text, chars_per_line):
    return max(1, -(-len(str(text)) // chars_per_line))


_FONT_FILES = {
    False: "/Users/docbrown/Library/Fonts/proximanova_regular.ttf",
    True: "/Users/docbrown/Library/Fonts/proximanova_bold.otf",
}
_font_cache = {}


def _wrap_lines(text, px_size, width_px, bold=False):
    """Точное число строк при переносе по словам — промер РЕАЛЬНЫМ шрифтом
    (оценка по числу символов врала → боксы наезжали друг на друга, 28.07)."""
    text = str(text or "").strip()
    if not text:
        return 0
    try:
        from PIL import ImageFont
        key = (bold, int(px_size * 4))
        font = _font_cache.get(key)
        if font is None:
            font = ImageFont.truetype(_FONT_FILES[bold], int(px_size * 4))
            _font_cache[key] = font
        scale = 4.0  # промер в 4× для точности на дробных кеглях
        lines, cur = 1, 0.0
        space_w = font.getlength(" ") / scale
        for word in text.split():
            w = font.getlength(word) / scale
            if cur > 0 and cur + space_w + w > width_px:
                lines += 1
                cur = w
            else:
                cur = cur + (space_w if cur > 0 else 0) + w
        return lines
    except Exception:
        return _est_lines(text, max(10, int(width_px / (px_size * 0.55))))


def build_pptx(data: dict, out_path: Path, prefs: dict = None) -> Path:
    from selection_pdf import _fmt_price, _cert_line
    prefs = prefs or {}
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Единый кегль текста для ВСЕЙ презентации — по самому длинному слайду
    # (по-слайдовое ужатие давало разнобой размеров; Вашик, 20.07).
    # База 15.5pt = 31px PDF; text_size="large" (34px) → +1.5pt как в PDF-префе.
    max_total = max((len((it.get("blurb") or "")) + len(it.get("headline") or "")
                     for it in data.get("items", [])), default=0)
    body_size = 15.5 if max_total <= 450 else 14 if max_total <= 700 else 13 if max_total <= 950 else 12
    if prefs.get("text_size") == "large":
        body_size += 1.5
    body_px = body_size * 2          # обратно в px холста для расчёта высот
    line_px = body_px * 1.31         # line-height 1.31 как в PDF

    with tempfile.TemporaryDirectory() as td:
        workdir = Path(td)
        for idx, item in enumerate(data.get("items", [])):
            slide = prs.slides.add_slide(blank)

            # ── фото-колонка: варианты A/B/C как в PDF ──
            main = _fetch(item.get("photo_product") or "", workdir, idx * 10)
            details_raw = [d for d in (item.get("photo_details") or []) if d][:4]
            details = [_fetch(d, workdir, idx * 10 + 1 + j) for j, d in enumerate(details_raw)]
            details = [d for d in details if d]
            if not details:
                # B: одно фото на всю колонку, рамка на самом фото
                _pic_contain_selfborder(slide, main, 60, 60, 853, 1320)
            elif len(details) == 1:
                # A: экспонат в рамке-контейнере + деталь cover-заливкой
                _pic_contain_framed(slide, main, 60, 60, 853, 640)
                _pic_cover(slide, details[0], 60, 740, 853, 640, workdir, f"{idx}d0")
            else:
                # C: экспонат + ряд/сетка квадратов деталей (как PDF .sig-row)
                _pic_contain_framed(slide, main, 60, 60, 853, 640)
                n = len(details)
                if n == 4:
                    side = 220
                    grid_w = 2 * side + 18
                    x0 = 60 + (853 - grid_w) / 2
                    y0 = 740 + (640 - grid_w) / 2
                    for j, d in enumerate(details):
                        cx = x0 + (j % 2) * (side + 18)
                        cy = y0 + (j // 2) * (side + 18)
                        _pic_cover(slide, d, cx, cy, side, side, workdir, f"{idx}d{j}")
                else:
                    side = min(620, (853 - 18 * (n - 1)) // n - 6)
                    row_w = n * side + 18 * (n - 1)
                    x0 = 60 + (853 - row_w) / 2
                    y0 = 740 + (640 - side) / 2
                    for j, d in enumerate(details):
                        _pic_cover(slide, d, x0 + j * (side + 18), y0, side, side,
                                   workdir, f"{idx}d{j}")

            # ── лого (как в PDF: центр текст-колонки; преф right — правый угол) ──
            logo = ASSETS / "logo_block.png"
            if logo.exists():
                try:
                    if prefs.get("logo_position") == "right":
                        slide.shapes.add_picture(str(logo), _in(1920 - 60 - 260), _in(70), width=_in(260))
                    else:
                        slide.shapes.add_picture(str(logo), _in(1251), _in(70), width=_in(330))
                except Exception:
                    pass

            # ── текстовая колонка: отдельные текстбоксы (фидбек Насти 28.07),
            #    отступы и веса — точно из PDF CSS ──
            col_x, col_w = _in(973), _in(887)
            y = 430.0  # px холста

            person = item.get("person") or ""
            n = _wrap_lines(person, 48, 887, bold=True)
            name_h = 48 * 1.3 * n
            tf = _txbox(slide, col_x, _in(y), col_w, _in(name_h + 10))
            p0 = tf.paragraphs[0]; p0.line_spacing = 1.3
            r = p0.add_run(); r.text = person
            _style(r, 24, bold=True, color=BLACK)
            y += name_h

            def _text_block(text, top_margin, size=None, bold=False, color=DARK):
                """Отдельный текстбокс с PDF-отступом; возвращает новую y."""
                nonlocal y
                size = size or body_size
                px = size * 2
                nlines = _wrap_lines(text, px, 887, bold=bold)
                h = px * 1.31 * nlines
                y_new = y + top_margin
                tfb = _txbox(slide, col_x, _in(y_new), col_w, _in(h + 8))
                par = tfb.paragraphs[0]; par.line_spacing = 1.31
                run = par.add_run(); run.text = text
                _style(run, size, bold=bold, color=color)
                y = y_new + h
                return tfb

            if item.get("headline"):
                _text_block(item["headline"], 34)                       # .subtitle

            price = _fmt_price(item.get("price"))
            price_disc = _fmt_price(item.get("price_discounted"))
            if price and price_disc:
                _text_block(f"Цена: {price}", 42)                       # .price-was
                _text_block(f"Цена с учётом вашей скидки: {price_disc}",
                            16, size=19.5, bold=True, color=BLACK)      # .price-discounted
            elif price:
                if prefs.get("price_bold"):
                    _text_block(price, 42, size=19.5, bold=True, color=BLACK)
                else:
                    _text_block(price, 42)                              # .price-line

            if item.get("dimensions"):
                _text_block(f"Размеры: {item['dimensions']}", 18)       # .meta

            cert = _cert_line(item.get("cert"))
            if cert:
                _text_block(cert, 42)                                   # .cert

            blurb = (item.get("blurb") or "").strip()
            paras = [x.strip() for x in blurb.split("\n\n") if x.strip()]
            if paras:                                                   # .blurb
                nlines = sum(_wrap_lines(x, body_px, 887) for x in paras)
                h = line_px * nlines + 20 * (len(paras) - 1)
                y_new = y + 42
                tfb = _txbox(slide, col_x, _in(y_new), col_w, _in(min(h + 10, 1380 - y_new)))
                first = True
                for para in paras:
                    p = tfb.paragraphs[0] if first else tfb.add_paragraph()
                    if not first:
                        p.space_before = Pt(20)  # 40px между абзацами, как в PDF
                    first = False
                    r = p.add_run(); r.text = para
                    _style(r, body_size)
                    p.line_spacing = 1.31
                y = y_new + min(h, 1380 - y_new)

            url = (item.get("url") or "").strip()
            if url and prefs.get("show_link", True):                    # .more
                # PDF-эталон: тёмная строка в потоке (+46px), подчёркнуто и
                # кликабельно только «Имя — Тип»
                link_name = person + (" — " + item["headline"] if item.get("headline") else "")
                full = "Больше информации об экспонате: " + link_name
                nlines = _wrap_lines(full, body_px, 887)
                y_new = min(y + 46, 1380 - line_px * nlines)
                link_tf = _txbox(slide, col_x, _in(y_new), col_w, _in(line_px * nlines + 8))
                lp = link_tf.paragraphs[0]; lp.line_spacing = 1.31
                pre = lp.add_run(); pre.text = "Больше информации об экспонате: "
                _style(pre, body_size)
                lr = lp.add_run(); lr.text = link_name
                _style(lr, body_size, underline=True)
                try:
                    lr.hyperlink.address = url
                except Exception:
                    pass

        # ── финальный слайд (зеркало PDF: колонка 90/500 отцентрована по
        #    вертикали, текст 17pt light, сетка 2×2 605×650 с зазором 20) ──
        try:
            if not prefs.get("show_final_slide", True):
                raise ValueError("final slide off")
            from selection_pdf import FINAL_SLIDE_TEXT
            texts = tuple(FINAL_SLIDE_TEXT)
            if (prefs.get("final_text") or "").strip():
                texts = tuple(t.strip() for t in prefs["final_text"].split("\n\n") if t.strip())
            slide = prs.slides.add_slide(blank)

            logo = ASSETS / "logo_block.png"
            logo_h = 100.0
            if logo.exists():
                try:
                    from PIL import Image
                    lw, lh = Image.open(logo).size
                    logo_h = 330 * lh / lw
                except Exception:
                    pass
            text_h = sum(_wrap_lines(t, 34, 500) * 34 * 1.4 for t in texts) + 46 * (len(texts) - 1)
            total_h = logo_h + 80 + text_h
            y0 = max(60.0, (1440 - total_h) / 2)

            if logo.exists():
                slide.shapes.add_picture(str(logo), _in(90), _in(y0), width=_in(330))
            tf = _txbox(slide, _in(90), _in(y0 + logo_h + 80), _in(500), _in(text_h + 20))
            first = True
            for t in texts:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                if not first:
                    p.space_before = Pt(23)  # 46px как в PDF
                first = False
                r = p.add_run(); r.text = t
                _style(r, 17)                # 34px как .final-text
                p.line_spacing = 1.4
            for i, n in enumerate((0, 1, 2, 3)):
                img = ASSETS / f"pg20-00{n}.jpg"
                if img.exists():
                    col, row = i % 2, i // 2
                    slide.shapes.add_picture(
                        str(img), _in(630 + col * 625), _in(60 + row * 670),
                        width=_in(605), height=_in(650))
        except Exception:
            pass

    prs.save(str(out_path))
    return out_path

def convert_to_key(pptx_path, key_path) -> bool:
    """PPTX → родной Keynote (.key) через Keynote.app (Вашик, 20.07: в .key ничего
    не съезжает и ссылки тёмные — Keynote-рендер и есть эталонный вид).
    Гоча: первый AppleScript-вызов после простоя таймаутится (-1712) — activate+retry."""
    import subprocess, os
    script = f"""
with timeout of 300 seconds
tell application "Keynote"
    activate
    set theDoc to open POSIX file "{pptx_path}"
    delay 3
    save theDoc in POSIX file "{key_path}"
    close theDoc saving no
end tell
end timeout
"""
    for attempt in (1, 2):
        try:
            r = subprocess.run(["osascript", "-e", script],
                               capture_output=True, text=True, timeout=330)
            if r.returncode == 0 and os.path.exists(str(key_path)):
                return True
        except Exception:
            pass
    return False
